"""
Generate an IEEE-style PowerPoint presentation for:
"Reinforcement Learning for Thrust-Vectored Fin Control of an EDF VTOL Vehicle in Simulation"

Targeted at a general AI conference audience, ~20 minutes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette (IEEE-inspired, professional) ──
DARK_BLUE = RGBColor(0x00, 0x3B, 0x6F)
MEDIUM_BLUE = RGBColor(0x00, 0x5A, 0x9C)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x00, 0x7A, 0x33)
RED = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── Helper functions ──

def add_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, items,
                              font_size=20, color=DARK_GRAY, bullet_color=MEDIUM_BLUE,
                              font_name="Calibri", spacing=8, line_spacing=1.15):
    """Add bulleted text content to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullets (tuples: (text, level))
        if isinstance(item, tuple):
            text, level = item
            indent = level
        else:
            text = item
            indent = 0

        bullet_char = "●" if indent == 0 else "○"
        prefix = "    " * indent
        p.text = f"{prefix}{bullet_char}  {text}"
        p.font.size = Pt(font_size - (indent * 2))
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(spacing)
        p.space_before = Pt(0)
        if line_spacing != 1.0:
            p.line_spacing = Pt((font_size - (indent * 2)) * line_spacing)

    return txBox


def add_section_header(slide, left, top, width, text, font_size=14,
                        color=MEDIUM_BLUE):
    """Add a small section header / label."""
    return add_textbox(slide, left, top, width, Inches(0.4), text,
                       font_size=font_size, color=color, bold=True,
                       font_name="Calibri")


def add_slide_number(slide, number, total=17):
    """Add slide number in bottom-right."""
    add_textbox(slide, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.35),
                f"{number}/{total}", font_size=11, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_footer_bar(slide):
    """Add a thin colored bar at the bottom."""
    add_rect(slide, Inches(0), Inches(7.25), SLIDE_WIDTH, Inches(0.25), DARK_BLUE)


def add_header_bar(slide, title_text, slide_num, total=17):
    """Add standard header bar with title and slide number."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.95), DARK_BLUE)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
                title_text, font_size=30, color=WHITE, bold=True)
    add_footer_bar(slide)
    add_slide_number(slide, slide_num, total)


def add_two_column_header(slide, title_text, slide_num, total=17):
    """Same header, just returns for convenience."""
    add_header_bar(slide, title_text, slide_num, total)


def make_table(slide, left, top, width, height, rows, cols, data,
               col_widths=None, header_color=DARK_BLUE, font_size=14):
    """Create a formatted table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_GRAY
                    paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_background(slide, WHITE)

# Top blue band
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(3.2), DARK_BLUE)
# Accent stripe
add_rect(slide, Inches(0), Inches(3.2), SLIDE_WIDTH, Inches(0.08), ACCENT_ORANGE)

# Title
add_textbox(slide, Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2),
            "Reinforcement Learning for Thrust-Vectored\nFin Control of an EDF VTOL Vehicle",
            font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT, line_spacing=1.15)

# Subtitle
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.8),
            "PPO-Based Attitude and Position Control in GPU-Accelerated Simulation",
            font_size=22, color=RGBColor(0xBB, 0xD5, 0xED), bold=False,
            alignment=PP_ALIGN.LEFT)

# Author info
add_textbox(slide, Inches(0.8), Inches(3.8), Inches(6), Inches(0.5),
            "Tang Zijian (Jacob Tang)", font_size=22, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(4.35), Inches(6), Inches(0.8),
            "Department of Applied Data Science\nMinnesota State University, Mankato",
            font_size=16, color=MEDIUM_GRAY)

# Conference / date placeholder
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
            "AI Conference 2026", font_size=16, color=MEDIUM_GRAY, bold=True)

# Image placeholder note (right side)
add_rect(slide, Inches(8.0), Inches(3.6), Inches(4.5), Inches(3.2), LIGHT_GRAY)
add_textbox(slide, Inches(8.2), Inches(4.6), Inches(4.1), Inches(1.0),
            "[Insert EDF Drone Photo / CAD Render]",
            font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide)
add_slide_number(slide, 1)



# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Outline / Agenda
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Presentation Outline", 2)

agenda_items = [
    ("1", "Motivation & Problem Statement", "Why thrust-vectored EDF drones?"),
    ("2", "Vehicle Platform & Control Architecture", "The physical testbed"),
    ("3", "Simulation Environment", "GPU-accelerated Isaac Sim setup"),
    ("4", "Physics Models", "Propulsion, aerodynamics, actuators"),
    ("5", "PPO Algorithm & Training", "Observations, actions, reward design"),
    ("6", "Experimental Results", "Hover & landing performance"),
    ("7", "Discussion & Future Work", "Challenges and next steps"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.3) + Inches(i * 0.78)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MEDIUM_BLUE if i < 7 else LIGHT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, Inches(1.8), y - Inches(0.02), Inches(5), Inches(0.35),
                title, font_size=20, color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(1.8), y + Inches(0.30), Inches(5), Inches(0.30),
                desc, font_size=14, color=MEDIUM_GRAY)

    # Connecting line (except last)
    if i < len(agenda_items) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(1.25), y + Inches(0.55),
                                       Inches(0.05), Inches(0.23))
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT_BLUE
        line.line.fill.background()

# Timing note
add_textbox(slide, Inches(8.5), Inches(1.5), Inches(4), Inches(0.4),
            "~20 minute presentation", font_size=16, color=ACCENT_ORANGE, bold=True)

# Key takeaway box
box = add_rect(slide, Inches(8.0), Inches(2.2), Inches(4.8), Inches(3.5), LIGHT_BLUE)
add_textbox(slide, Inches(8.3), Inches(2.4), Inches(4.2), Inches(0.4),
            "Key Takeaway", font_size=18, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(8.3), Inches(2.9), Inches(4.2), Inches(2.5),
            "Can deep RL learn to control a vehicle that steers by deflecting fins in its own exhaust stream?\n\n"
            "Yes — PPO achieves 100% landing rate with soft touchdowns in simulation, "
            "opening the door to sim-to-real transfer on a physical drone testbed.",
            font_size=15, color=DARK_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Motivation & Problem Statement
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Motivation: Why Thrust-Vectored EDF Drones?", 3)

# Left column
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(5.5),
                   "THE CONTROL CHALLENGE")
add_bullet_slide_content(slide, Inches(0.6), Inches(1.55), Inches(5.8), Inches(5.0), [
    "Conventional multirotors: independent rotor speeds → decoupled control",
    "Thrust-vectored EDF: single thrust source + 4 jet-vane fins",
    "Fin effectiveness ∝ throttle² — at 50% throttle, only 25% authority",
    "Nonlinear aerodynamics, actuator lag, gyroscopic precession",
    "Classical linear controllers (PID/LQR) limited to narrow operating envelope",
    "RL can learn the full nonlinear mapping end-to-end",
])

# Right column — "Why it matters" box
add_rect(slide, Inches(7.0), Inches(1.15), Inches(5.8), Inches(2.8), LIGHT_BLUE)
add_textbox(slide, Inches(7.3), Inches(1.3), Inches(5.2), Inches(0.4),
            "Why It Matters", font_size=18, color=DARK_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(7.3), Inches(1.75), Inches(5.2), Inches(2.0), [
    "Aerodynamic efficiency over multirotors",
    "Noise reduction (single fan vs. 4+ rotors)",
    "Mechanical simplicity for certain missions",
    "Applicable to rocket landing (jet vanes)",
], font_size=16, spacing=6)

# Key equation box
add_rect(slide, Inches(7.0), Inches(4.3), Inches(5.8), Inches(1.8), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(7.3), Inches(4.45), Inches(5.2), Inches(0.35),
            "The Core Coupling", font_size=16, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(7.3), Inches(4.85), Inches(5.2), Inches(1.0),
            "F_fin = q · S · C_N(α)    where    q ∝ throttle²\n\n"
            "→ Reducing throttle to descend simultaneously\n"
            "   reduces your ability to steer laterally",
            font_size=15, color=DARK_GRAY, line_spacing=1.3)



# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Related Work
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Related Work & Research Gap", 4)

# Three columns
col_data = [
    ("RL for Aerial Vehicles", [
        "Hwangbo et al. — PPO for agile quadrotor control",
        "Koch et al. — Deep RL for attitude control",
        "Pi et al. — Sim-to-real under wind disturbances",
        "GPU-parallel training (Isaac Gym, Rudin et al.)",
    ]),
    ("Thrust Vector Control", [
        "Long history in rocketry (gimbaled nozzles)",
        "Jet-vane TVC for VTOL vehicles",
        "Classical: PID/LQR around hover linearization",
        "Limited operating envelope for nonlinear regime",
    ]),
    ("Sim-to-Real Transfer", [
        "Domain randomization (Tobin et al.)",
        "System identification (Yu et al.)",
        "CAD2RL: simulation-only to real flight",
        "Fidelity of dynamics model is critical",
    ]),
]

for i, (title, items) in enumerate(col_data):
    x = Inches(0.5) + Inches(i * 4.2)
    add_rect(slide, x, Inches(1.15), Inches(3.9), Inches(0.5), MEDIUM_BLUE)
    add_textbox(slide, x + Inches(0.15), Inches(1.2), Inches(3.6), Inches(0.4),
                title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_slide_content(slide, x + Inches(0.15), Inches(1.8), Inches(3.6), Inches(3.5),
                              items, font_size=14, spacing=6, line_spacing=1.2)

# Gap callout
add_rect(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.3), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(0.8), Inches(5.6), Inches(2.0), Inches(0.4),
            "RESEARCH GAP", font_size=16, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
            "RL for thrust-vectored fin control of EDF vehicles remains largely unexplored. "
            "Key distinctions: single thrust source with coupled fins, quadratic throttle-authority dependence, "
            "hobby-grade actuator dynamics, and gyroscopic effects from high-speed rotor.",
            font_size=15, color=DARK_GRAY, line_spacing=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Vehicle Platform
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Vehicle Platform: EDF TVC Drone Testbed", 5)

# Left: specs
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(5.5),
                   "PHYSICAL SPECIFICATIONS")

spec_data = [
    ["Parameter", "Value"],
    ["Total mass", "3.1 kg"],
    ["Body length / diameter", "0.35 m / 0.12 m"],
    ["Ixx, Iyy / Izz", "0.05 / 0.02 kg·m²"],
    ["EDF unit", "FMS 90mm 12-blade"],
    ["Motor", "4068-KV1850 inrunner"],
    ["Battery", "6S LiPo"],
    ["Servos", "4× MG996R"],
    ["Fin count / area", "4 / 0.002 m² each"],
    ["Max fin deflection", "±15° (0.262 rad)"],
    ["Max thrust", "39.2 N"],
]

make_table(slide, Inches(0.6), Inches(1.55), Inches(5.5), Inches(4.8),
           len(spec_data), 2, spec_data,
           col_widths=[Inches(3.2), Inches(2.3)], font_size=13)

# Right: image placeholder + control architecture
add_rect(slide, Inches(7.0), Inches(1.15), Inches(5.8), Inches(2.8), LIGHT_GRAY)
add_textbox(slide, Inches(7.5), Inches(2.2), Inches(4.8), Inches(0.5),
            "[Insert CAD Render: cad_v1.png]",
            font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Control architecture box
add_rect(slide, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.6), LIGHT_BLUE)
add_textbox(slide, Inches(7.3), Inches(4.35), Inches(5.2), Inches(0.35),
            "Control Architecture", font_size=16, color=DARK_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(7.3), Inches(4.75), Inches(5.2), Inches(2.0), [
    "Forward fin (+X): pitch control",
    "Right fin (+Y): roll control",
    "Aft fin (−X): complementary pitch",
    "Left fin (−Y): complementary roll",
    "Yaw: differential tangential forces",
    "Altitude: EDF throttle",
    "→ 5 actuators → full 6-DOF control",
], font_size=14, spacing=4)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Simulation Environment Overview
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Simulation Environment: NVIDIA Isaac Sim", 6)

# Left column
add_bullet_slide_content(slide, Inches(0.6), Inches(1.25), Inches(5.8), Inches(5.5), [
    "Built on Isaac Lab v2.3.2 (PhysX 5 GPU physics)",
    "128 parallel environments on a single GPU",
    "Gymnasium-compatible vectorized interface",
    "Physics timestep: 8.33 ms (120 Hz)",
    "RL decision timestep: 33.3 ms (30 Hz, decimation = 4)",
    "Semi-empirical dynamics models calibrated to hardware",
    "Two force dispatch modes:",
    ("Per-link: forces at center-of-pressure locations", 1),
    ("Collapsed body wrench: for debugging/validation", 1),
], font_size=18, spacing=6)

# Right: image placeholder
add_rect(slide, Inches(7.0), Inches(1.15), Inches(5.8), Inches(3.5), LIGHT_GRAY)
add_textbox(slide, Inches(7.5), Inches(2.5), Inches(4.8), Inches(0.5),
            "[Insert: isaacsim_128_env.png\n128 parallel environments]",
            font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Key numbers box
add_rect(slide, Inches(7.0), Inches(4.9), Inches(5.8), Inches(1.8), LIGHT_BLUE)
add_textbox(slide, Inches(7.3), Inches(5.0), Inches(5.2), Inches(0.35),
            "Key Numbers", font_size=16, color=DARK_BLUE, bold=True)

nums = [("128", "Parallel Envs"), ("120 Hz", "Physics Rate"),
        ("30 Hz", "RL Decision Rate"), ("4×", "Decimation")]
for i, (val, label) in enumerate(nums):
    x = Inches(7.3) + Inches(i * 1.3)
    add_textbox(slide, x, Inches(5.45), Inches(1.2), Inches(0.5),
                val, font_size=22, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(5.95), Inches(1.2), Inches(0.4),
                label, font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)



# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Propulsion Model
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Propulsion Model: EDF Dynamics", 7)

# Equations (left)
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(5.5),
                   "FIRST-ORDER SPOOL + QUADRATIC THRUST")

add_textbox(slide, Inches(0.8), Inches(1.65), Inches(6.0), Inches(0.5),
            "ω̇ = (ω_target − ω) / τ_motor        (spool dynamics)",
            font_size=20, color=DARK_GRAY, font_name="Consolas")
add_textbox(slide, Inches(0.8), Inches(2.25), Inches(6.0), Inches(0.5),
            "T = k_T · ω²                          (quadratic thrust)",
            font_size=20, color=DARK_GRAY, font_name="Consolas")

add_section_header(slide, Inches(0.6), Inches(3.0), Inches(5.5),
                   "THREE TORQUE COMPONENTS")
add_bullet_slide_content(slide, Inches(0.6), Inches(3.4), Inches(6.0), Inches(3.0), [
    "Static reaction torque: τ = −k_Q · ω² (opposes rotor spin)",
    "Dynamic spool torque: τ = −I_rotor · ω̇ (reaction to acceleration)",
    "Gyroscopic precession: τ = −ω_body × (I_rotor · ω · ê_z)",
    ("Couples body rotation to rotor angular momentum", 1),
    ("Critical for realistic yaw dynamics", 1),
], font_size=16, spacing=5)

# Right: parameters table
add_section_header(slide, Inches(7.0), Inches(1.15), Inches(5.5),
                   "EDF PARAMETERS")

edf_data = [
    ["Parameter", "Value"],
    ["Max thrust T_max", "39.2 N"],
    ["Motor time constant τ", "0.15 s"],
    ["Max angular velocity ω_max", "4300 rad/s"],
    ["Rotor inertia I_rotor", "0.0002 kg·m²"],
    ["Exhaust speed (nominal)", "116 m/s"],
    ["Gyro torque scale", "0.1"],
]

make_table(slide, Inches(7.0), Inches(1.55), Inches(5.5), Inches(3.0),
           len(edf_data), 2, edf_data,
           col_widths=[Inches(3.2), Inches(2.3)], font_size=13)

# Insight box
add_rect(slide, Inches(7.0), Inches(4.9), Inches(5.8), Inches(1.8), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(7.3), Inches(5.0), Inches(5.2), Inches(0.35),
            "Key Insight", font_size=16, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(7.3), Inches(5.4), Inches(5.2), Inches(1.0),
            "The gyroscopic precession torque couples body rotation to rotor angular momentum — "
            "a real physical effect that makes yaw control non-trivial and distinguishes this "
            "from simplified point-mass models.",
            font_size=14, color=DARK_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Fin Aerodynamic Model
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Fin Aerodynamic Model: Jet-Vane Forces", 8)

# Left: equations
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(6.0),
                   "SEMI-EMPIRICAL FORCE MODEL")

eqs = [
    ("Dynamic pressure:", "q = ½ρ · v²_exhaust · throttle² · k_duct"),
    ("Normal force coeff:", "C_N(α) = C_Nα · α · (1 − k_sat · α²)"),
    ("Drag coeff:", "C_D(α) = C_D0 + C_Dα² · α²"),
    ("Per-fin forces:", "F_n = q · S_fin · C_N(α),  F_t = q · S_fin · C_D(α)"),
]

for i, (label, eq) in enumerate(eqs):
    y = Inches(1.6) + Inches(i * 0.75)
    add_textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.35),
                label, font_size=14, color=MEDIUM_BLUE, bold=True)
    add_textbox(slide, Inches(0.8), y + Inches(0.3), Inches(6.0), Inches(0.35),
                eq, font_size=17, color=DARK_GRAY, font_name="Consolas")

# Right: parameters
aero_data = [
    ["Parameter", "Value"],
    ["Air density ρ", "1.225 kg/m³"],
    ["Exhaust speed v_exh", "116 m/s"],
    ["Duct correction k_duct", "1.3"],
    ["Lift slope C_Nα", "3.5 rad⁻¹"],
    ["Stall saturation k_sat", "2.0"],
    ["Zero-defl drag C_D0", "0.05"],
    ["Drag slope C_Dα²", "1.5"],
    ["Fin area S_fin", "0.002 m²"],
]

make_table(slide, Inches(7.0), Inches(1.15), Inches(5.5), Inches(4.0),
           len(aero_data), 2, aero_data,
           col_widths=[Inches(3.2), Inches(2.3)], font_size=13)

# Critical feature callout
add_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.5), RGBColor(0xFC, 0xE4, 0xEC))
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.35),
            "CRITICAL COUPLING", font_size=16, color=RED, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
            "Fin force ∝ throttle² — at 50% throttle, fin forces drop to 25% of maximum. "
            "This quadratic coupling between thrust and control authority is THE defining challenge "
            "of jet-vane TVC and why classical linear controllers struggle outside the hover envelope.",
            font_size=16, color=DARK_GRAY, line_spacing=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — PPO Algorithm
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Learning Algorithm: Proximal Policy Optimization", 9)

# Left: PPO overview
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(6.0),
                   "PPO WITH CLIPPED SURROGATE OBJECTIVE")

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6.0), Inches(0.5),
            "L = E[ min( r(θ)·Â, clip(r(θ), 1−ε, 1+ε)·Â ) ]",
            font_size=20, color=DARK_GRAY, font_name="Consolas")

add_bullet_slide_content(slide, Inches(0.6), Inches(2.3), Inches(6.0), Inches(4.0), [
    "On-policy actor-critic algorithm",
    "Trust-region constraint prevents destructive updates",
    "Clip coefficient ε = 0.2",
    "GAE advantage estimation (γ=0.99, λ=0.95)",
    "Adaptive KL early stopping (KL_target = 0.03)",
    "Combined loss: policy + value (0.5) + entropy (0.005)",
], font_size=17, spacing=6)

# Right: hyperparameters table
add_section_header(slide, Inches(7.0), Inches(1.15), Inches(5.5),
                   "HYPERPARAMETERS (LANDING TASK)")

hyper_data = [
    ["Parameter", "Value"],
    ["Parallel environments", "128"],
    ["Rollout length", "128 steps"],
    ["Minibatches / update", "8"],
    ["Update epochs", "4"],
    ["Learning rate", "3 × 10⁻⁴"],
    ["Discount γ", "0.99"],
    ["GAE λ", "0.95"],
    ["Clip ε", "0.2"],
    ["Entropy coeff", "0.005"],
    ["Max grad norm", "0.5"],
]

make_table(slide, Inches(7.0), Inches(1.55), Inches(5.5), Inches(4.8),
           len(hyper_data), 2, hyper_data,
           col_widths=[Inches(3.2), Inches(2.3)], font_size=13)

# Why PPO box
add_rect(slide, Inches(0.5), Inches(5.8), Inches(5.8), Inches(1.0), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(5.9), Inches(5.2), Inches(0.7),
            "Why PPO? Robust training stability for continuous control, "
            "proven track record in aerial vehicle RL (Koch et al., Hwangbo et al.)",
            font_size=14, color=DARK_BLUE, line_spacing=1.25)



# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Observation & Action Space
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Observation & Action Space Design", 10)

# Left: observation space
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(5.5),
                   "24-DIMENSIONAL OBSERVATION VECTOR")

obs_data = [
    ["Index", "Component", "Dim"],
    ["0–2", "Position error (m)", "3"],
    ["3–6", "Attitude quaternion [w,x,y,z]", "4"],
    ["7–9", "Linear velocity, body FRD", "3"],
    ["10–12", "Angular velocity, body FRD", "3"],
    ["13", "Height above ground", "1"],
    ["14–17", "Fin angles (rad)", "4"],
    ["18–21", "Fin angular rates (rad/s)", "4"],
    ["22", "Motor RPM, normalized", "1"],
    ["23", "Contact state", "1"],
]

make_table(slide, Inches(0.5), Inches(1.55), Inches(6.2), Inches(4.2),
           len(obs_data), 3, obs_data,
           col_widths=[Inches(1.0), Inches(3.8), Inches(1.4)], font_size=13)

# Right: action space
add_section_header(slide, Inches(7.0), Inches(1.15), Inches(5.5),
                   "5-DIMENSIONAL ACTION VECTOR")

add_textbox(slide, Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.5),
            "a = [α₁, α₂, α₃, α₄, u_throttle]",
            font_size=20, color=DARK_GRAY, font_name="Consolas")

add_bullet_slide_content(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(2.5), [
    "4 fin deflections: ±15° (±0.262 rad)",
    "1 throttle command: [0, 1]",
    "Network outputs tanh → linear scaling",
    "Smaller than quadrotor (typically 4 rotors)",
], font_size=16, spacing=6)

# Network architecture box
add_rect(slide, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.6), LIGHT_BLUE)
add_textbox(slide, Inches(7.3), Inches(4.35), Inches(5.2), Inches(0.35),
            "Network Architecture", font_size=16, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(7.3), Inches(4.8), Inches(5.2), Inches(0.4),
            "Actor:  Obs(24) → FC(256) → tanh → FC(256) → tanh → FC(5)",
            font_size=13, color=DARK_GRAY, font_name="Consolas")
add_textbox(slide, Inches(7.3), Inches(5.2), Inches(5.2), Inches(0.4),
            "Critic: Obs(24) → FC(256) → tanh → FC(256) → tanh → FC(1)",
            font_size=13, color=DARK_GRAY, font_name="Consolas")
add_textbox(slide, Inches(7.3), Inches(5.7), Inches(5.2), Inches(0.8),
            "Separate actor-critic networks\n"
            "Learnable per-dimension log-σ for exploration",
            font_size=14, color=MEDIUM_GRAY, line_spacing=1.3)

# Design rationale note
add_textbox(slide, Inches(0.6), Inches(6.0), Inches(6.0), Inches(0.8),
            "Design rationale: Quaternion avoids gimbal lock; fin angles/rates give "
            "proprioceptive feedback on actuator state (critical given servo lag); "
            "normalized RPM informs policy of available fin authority.",
            font_size=13, color=MEDIUM_GRAY, line_spacing=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Initialization Priors (Critical Insight)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Critical: Initialization Priors on Action Distribution", 11)

# Left: throttle bias
add_rect(slide, Inches(0.4), Inches(1.15), Inches(6.0), Inches(2.8), LIGHT_BLUE)
add_textbox(slide, Inches(0.7), Inches(1.25), Inches(5.4), Inches(0.35),
            "1. Throttle Bias Initialization", font_size=20, color=DARK_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(0.7), Inches(1.7), Inches(5.4), Inches(2.0), [
    "Output bias → atanh(2 × 0.78 − 1) ≈ 1.05",
    "Initial mean throttle ≈ 0.78 (near hover ≈ 0.88)",
    "Without this: default 0.5 → immediate free-fall",
    "Agent converges to zero-throttle local optimum",
    "Physical reasoning: √(mg / T_max) ≈ 0.88",
], font_size=16, spacing=5)

# Right: exploration noise
add_rect(slide, Inches(6.8), Inches(1.15), Inches(6.0), Inches(2.8), LIGHT_BLUE)
add_textbox(slide, Inches(7.1), Inches(1.25), Inches(5.4), Inches(0.35),
            "2. Per-Channel Exploration Noise", font_size=20, color=DARK_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(7.1), Inches(1.7), Inches(5.4), Inches(2.0), [
    "Fin channels: log σ = −2.0 (std ≈ 0.14)",
    "Throttle channel: log σ = −1.0 (std ≈ 0.37)",
    "Asymmetric: quieter fins, louder throttle",
    "Prevents lateral oscillations from fin noise",
    "Maintains altitude exploration range",
], font_size=16, spacing=5)

# Failure mode illustration
add_rect(slide, Inches(0.4), Inches(4.3), Inches(12.4), Inches(2.5), RGBColor(0xFC, 0xE4, 0xEC))
add_textbox(slide, Inches(0.7), Inches(4.4), Inches(11.8), Inches(0.35),
            "WITHOUT THESE PRIORS — FAILURE MODES", font_size=18, color=RED, bold=True)

add_textbox(slide, Inches(0.7), Inches(4.9), Inches(5.5), Inches(1.5),
            "Default throttle (0.5):\n"
            "→ Immediate free-fall\n"
            "→ Large negative per-step rewards\n"
            "→ Converges to zero-throttle equilibrium\n"
            "→ Agent learns: \"crash fast = minimize cost\"",
            font_size=15, color=DARK_GRAY, line_spacing=1.3)

add_textbox(slide, Inches(6.8), Inches(4.9), Inches(5.5), Inches(1.5),
            "Uniform exploration noise:\n"
            "→ Excessive lateral oscillations from fins\n"
            "→ Vehicle destabilized before learning descent\n"
            "→ No stable trajectories to learn from\n"
            "→ Policy never discovers productive control basin",
            font_size=15, color=DARK_GRAY, line_spacing=1.3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Reward Design
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Reward Design: Terminal vs. Per-Step Balance", 12)

# Left: reward table
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(5.5),
                   "LANDING TASK REWARD WEIGHTS")

reward_data = [
    ["Term", "Weight"],
    ["Alive bonus", "+0.10"],
    ["3D position error", "−0.05"],
    ["Horizontal position error", "−0.30"],
    ["Attitude error (tilt)", "−0.20"],
    ["Angular velocity", "−0.10"],
    ["Control effort (fin |α|)", "−0.02"],
    ["Control rate (fin |α̇|)", "−0.01"],
    ["Vertical speed shaping", "−0.30"],
    ["Delta-v cost (T/T_max)²", "−0.20"],
    ["─── Terminal ───", "───"],
    ["Crash penalty", "−200"],
    ["Touchdown softness", "+25·e^(−v)"],
    ["Pad accuracy", "+200·e^(−d)"],
    ["Landing success", "+250"],
]

make_table(slide, Inches(0.5), Inches(1.55), Inches(5.5), Inches(5.2),
           len(reward_data), 2, reward_data,
           col_widths=[Inches(3.5), Inches(2.0)], font_size=12)

# Right: key insight
add_rect(slide, Inches(6.8), Inches(1.15), Inches(6.0), Inches(3.0), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(7.1), Inches(1.25), Inches(5.4), Inches(0.35),
            "THE MAGNITUDE BALANCE PROBLEM", font_size=18, color=ACCENT_ORANGE, bold=True)
add_textbox(slide, Inches(7.1), Inches(1.75), Inches(5.4), Inches(2.2),
            "For episode length T = 900 (30s at 30Hz):\n\n"
            "Integrated per-step costs: Σ|w_i|·φ̄_i·T ≈ O(100)\n\n"
            "If terminal rewards < O(100):\n"
            "→ PPO optimizes per-step minimization\n"
            "→ Zero throttle = crash fast = minimize cost\n\n"
            "Solution: Terminal magnitudes (≥200) must\n"
            "clearly dominate integrated per-step budget",
            font_size=15, color=DARK_GRAY, line_spacing=1.3)

# Horizontal guidance
add_rect(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.2), LIGHT_BLUE)
add_textbox(slide, Inches(7.1), Inches(4.6), Inches(5.4), Inches(0.35),
            "Horizontal Guidance Shaping", font_size=16, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(7.1), Inches(5.0), Inches(5.4), Inches(1.5),
            "● Dense lateral gradient (w = −0.30) throughout descent\n"
            "● Independent of altitude — guides even at high altitudes\n"
            "● Exponential pad bonus e^(−d) has negligible gradient for d > 2m\n"
            "● Vertical speed shaping targets ~0.5 m/s descent rate",
            font_size=14, color=DARK_GRAY, line_spacing=1.35)



# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Training Procedure & Curriculum
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Training Procedure & Spawn-Position Curriculum", 13)

# Left: two tasks
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(5.5),
                   "TWO TRAINING TASKS")

# Hover box
add_rect(slide, Inches(0.5), Inches(1.55), Inches(5.8), Inches(2.0), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.65), Inches(5.2), Inches(0.35),
            "Hover Stabilization", font_size=18, color=DARK_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(0.8), Inches(2.05), Inches(5.2), Inches(1.3), [
    "Residual-PID: PID baseline + PPO correction (scale=0.05)",
    "2,048 environment steps, LR = 10⁻⁶",
    "Rollout: 64 steps, 2 minibatches, 2 epochs",
], font_size=14, spacing=4)

# Landing box
add_rect(slide, Inches(0.5), Inches(3.8), Inches(5.8), Inches(2.0), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Inches(0.8), Inches(3.9), Inches(5.2), Inches(0.35),
            "Autonomous Landing", font_size=18, color=GREEN, bold=True)
add_bullet_slide_content(slide, Inches(0.8), Inches(4.3), Inches(5.2), Inches(1.3), [
    "Pure PPO (no PID baseline) — full trajectory learning",
    "5,000,000 environment steps, LR = 3×10⁻⁴",
    "Rollout: 128 steps, 8 minibatches, 4 epochs",
], font_size=14, spacing=4)

# Right: curriculum
add_section_header(slide, Inches(7.0), Inches(1.15), Inches(5.5),
                   "SPAWN-POSITION CURRICULUM")

add_rect(slide, Inches(6.8), Inches(1.55), Inches(6.0), Inches(5.0), LIGHT_GRAY)

# Curriculum stages
stages = [
    ("Start (0 steps)", "±0.5 m XY spawn box", "Learn on-pad touchdown"),
    ("Linear anneal", "±0.5 m → ±2.0 m", "Gradually increase lateral offset"),
    ("End (3M steps)", "±2.0 m XY (full range)", "Full task difficulty"),
    ("Altitude", "8–12 m (constant)", "Same throughout training"),
    ("Evaluation", "Full ±2.0 m range", "Always un-curricularized"),
]

for i, (stage, value, desc) in enumerate(stages):
    y = Inches(1.75) + Inches(i * 0.9)
    add_textbox(slide, Inches(7.1), y, Inches(2.5), Inches(0.35),
                stage, font_size=15, color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(9.5), y, Inches(3.0), Inches(0.35),
                value, font_size=15, color=DARK_GRAY, font_name="Consolas")
    add_textbox(slide, Inches(7.1), y + Inches(0.35), Inches(5.5), Inches(0.35),
                desc, font_size=12, color=MEDIUM_GRAY)

# Rationale
add_textbox(slide, Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.6),
            "Rationale: Starting with a narrow spawn box lets the policy learn soft-touchdown behavior first, "
            "then generalize to larger lateral offsets. Evaluation always uses the full range.",
            font_size=14, color=MEDIUM_GRAY, line_spacing=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Hover Results
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Results: Hover Stabilization", 14)

# Results table
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(5.5),
                   "HOVER PERFORMANCE METRICS")

hover_data = [
    ["Metric", "Value", "Threshold"],
    ["Mean position error", "0.085 m", "< 0.5 m ✓"],
    ["Max position error", "0.215 m", "< 0.5 m ✓"],
    ["Mean tilt", "< 0.001 rad", "—"],
    ["Max tilt", "< 0.001 rad", "—"],
    ["Mean angular rate", "1.1×10⁻⁵ rad/s", "—"],
    ["Evaluation", "PASSED", "—"],
]

make_table(slide, Inches(0.5), Inches(1.55), Inches(6.0), Inches(3.2),
           len(hover_data), 3, hover_data,
           col_widths=[Inches(2.5), Inches(2.0), Inches(1.5)], font_size=14)

# Key findings
add_rect(slide, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.7), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(5.1), Inches(5.4), Inches(0.35),
            "Key Findings", font_size=16, color=DARK_BLUE, bold=True)
add_bullet_slide_content(slide, Inches(0.8), Inches(5.5), Inches(5.4), Inches(1.0), [
    "PID provides effective attitude stabilization",
    "RL residual corrects for position drift",
    "Only 2,048 steps needed — very fast convergence",
], font_size=14, spacing=4)

# Right: big numbers
for i, (val, label, clr) in enumerate([
    ("0.085 m", "Mean Position Error", GREEN),
    ("< 0.001 rad", "Mean Tilt", GREEN),
    ("2,048", "Training Steps", MEDIUM_BLUE),
]):
    y = Inches(1.3) + Inches(i * 1.8)
    add_rect(slide, Inches(7.0), y, Inches(5.5), Inches(1.5), LIGHT_GRAY)
    add_textbox(slide, Inches(7.3), y + Inches(0.15), Inches(4.9), Inches(0.7),
                val, font_size=40, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.3), y + Inches(0.9), Inches(4.9), Inches(0.4),
                label, font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Landing Results
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Results: Autonomous Landing", 15)

# Left: results table
add_section_header(slide, Inches(0.6), Inches(1.15), Inches(6.0),
                   "LANDING EVALUATION METRICS (128 PARALLEL EPISODES)")

landing_data = [
    ["Metric", "Value"],
    ["Landing rate", "100%"],
    ["Mean touchdown speed", "0.126 m/s"],
    ["Max touchdown speed", "0.347 m/s"],
    ["Mean pad distance", "1.89 m"],
    ["Max pad distance", "4.56 m"],
    ["Mean throttle (eval)", "0.759"],
    ["Max downward speed", "9.75 m/s"],
    ["On-pad fraction (d < 0.5m)", "2%"],
]

make_table(slide, Inches(0.5), Inches(1.55), Inches(5.5), Inches(3.8),
           len(landing_data), 2, landing_data,
           col_widths=[Inches(3.5), Inches(2.0)], font_size=14)

# Right: big result numbers
results_big = [
    ("100%", "Landing Rate", GREEN),
    ("0.126 m/s", "Mean Touchdown Speed", GREEN),
    ("1.89 m", "Mean Pad Distance", ACCENT_ORANGE),
    ("5M", "Training Steps", MEDIUM_BLUE),
]

for i, (val, label, clr) in enumerate(results_big):
    x = Inches(7.0) + Inches((i % 2) * 3.0)
    y = Inches(1.3) + Inches((i // 2) * 2.0)
    add_rect(slide, x, y, Inches(2.7), Inches(1.7), LIGHT_GRAY)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.15), Inches(2.5), Inches(0.7),
                val, font_size=32, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(1.0), Inches(2.5), Inches(0.5),
                label, font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Assessment box
add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.9),
            "✓ 100% landing rate with soft touchdowns — the policy reliably descends and lands safely\n"
            "⚠ Lateral accuracy remains limited: 1.89 m mean pad distance, only 2% on-pad — "
            "the policy has not fully solved lateral guidance to the pad center",
            font_size=16, color=DARK_GRAY, line_spacing=1.35)



# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — Training Curves & Analysis
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Training Progression & Analysis", 16)

# Left: chart placeholder
add_rect(slide, Inches(0.5), Inches(1.15), Inches(6.5), Inches(4.0), LIGHT_GRAY)
add_textbox(slide, Inches(1.5), Inches(2.7), Inches(4.5), Inches(0.5),
            "[Insert: training_curves.png\nMean Episode Reward over 5M Steps]",
            font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Right: analysis
add_section_header(slide, Inches(7.5), Inches(1.15), Inches(5.5),
                   "TRAINING DYNAMICS")

add_bullet_slide_content(slide, Inches(7.5), Inches(1.55), Inches(5.3), Inches(4.0), [
    "0–1M steps: reward ≈ −3 (per-step costs dominate)",
    "1–2M steps: transition to positive rewards",
    "2–5M steps: steady improvement to ≈ +1.75",
    "Monotonic upward trend confirms learning",
    "100% landing rate achieved early, maintained throughout",
], font_size=16, spacing=8)

# Interpretation box
add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.35),
            "Interpreting the Reward Crossover", font_size=16, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.7),
            "Negative → Positive crossover at ~2M steps: integrated per-step costs (attitude, control effort, delta-v) "
            "give way to terminal landing bonuses (+250 success + up to +200 pad accuracy). "
            "This marks when the majority of episodes end in successful landings rather than timeouts or crashes.",
            font_size=15, color=DARK_GRAY, line_spacing=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — Discussion: Lateral Accuracy Challenge
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Discussion: The Lateral Accuracy Challenge", 17)

challenges = [
    ("Throttle–Fin Coupling",
     "Lateral authority ∝ throttle². During descent, reducing throttle simultaneously "
     "reduces fin effectiveness — fundamental tension between descent rate and lateral controllability."),
    ("Actuator Bandwidth",
     "MG996R servos (τ = 0.05s, 7.54 rad/s max rate) limit correction speed, "
     "especially during final approach where rapid corrections are needed."),
    ("Reward Gradient at Distance",
     "Exponential pad bonus e^(−d) has negligible gradient for d > 2m. "
     "Linear horizontal penalty (w = −0.30) may be insufficient vs. vertical/delta-v terms."),
    ("Curriculum Interaction",
     "Narrow initial spawn helps learn touchdown but may not provide sufficient "
     "training signal for large lateral corrections from the full ±2.0m range."),
]

for i, (title, desc) in enumerate(challenges):
    y = Inches(1.2) + Inches(i * 1.45)
    # Number
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05),
                                        Inches(0.45), Inches(0.45))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = ACCENT_ORANGE
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, Inches(1.3), y, Inches(11.5), Inches(0.35),
                title, font_size=18, color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(1.3), y + Inches(0.4), Inches(11.5), Inches(0.8),
                desc, font_size=14, color=DARK_GRAY, line_spacing=1.25)

# Comparison note
add_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
            "vs. Quadrotors: Independent, throttle-decoupled attitude authority. TVC's coupled control "
            "means standard quadrotor RL techniques don't directly transfer.",
            font_size=14, color=DARK_BLUE, line_spacing=1.2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — Sim-to-Real Considerations
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Sim-to-Real Transfer Considerations", 18)

# Gap sources
gaps = [
    ("Aerodynamic Model Fidelity",
     "Flat-plate fin model may miss flow separation, wake interactions between fins, "
     "and ground-effect changes during landing approach.",
     "CFD validation, wind tunnel data"),
    ("Servo Nonlinearities",
     "Real MG996R servos exhibit backlash, load-dependent speed variation, "
     "and temperature drift not captured by first-order model.",
     "System identification on hardware"),
    ("Battery Voltage Sag",
     "6S LiPo voltage drops under load, reducing available thrust during "
     "high-throttle maneuvers — not currently modeled.",
     "Voltage-dependent thrust model"),
    ("Sensor Noise",
     "Simulation provides clean state observations; physical testbed uses "
     "noisy IMU and position estimates.",
     "Observation noise injection during training"),
]

for i, (title, desc, mitigation) in enumerate(gaps):
    y = Inches(1.2) + Inches(i * 1.4)
    add_rect(slide, Inches(0.5), y, Inches(8.0), Inches(1.15),
             LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(7.5), Inches(0.3),
                title, font_size=16, color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(0.8), y + Inches(0.35), Inches(7.5), Inches(0.6),
                desc, font_size=13, color=DARK_GRAY, line_spacing=1.2)

    # Mitigation
    add_rect(slide, Inches(8.8), y, Inches(4.0), Inches(1.15),
             RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, Inches(9.0), y + Inches(0.05), Inches(3.6), Inches(0.25),
                "Mitigation:", font_size=12, color=GREEN, bold=True)
    add_textbox(slide, Inches(9.0), y + Inches(0.35), Inches(3.6), Inches(0.6),
                mitigation, font_size=13, color=DARK_GRAY, line_spacing=1.2)

# Domain randomization callout
add_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.5),
            "Primary strategy: Domain randomization of aerodynamic coefficients, servo parameters, "
            "mass properties, and sensor noise during training to improve transfer robustness.",
            font_size=15, color=DARK_BLUE, bold=False, line_spacing=1.2)



# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — Future Work
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Future Work", 19)

future_items = [
    ("GTrXL Policy",
     "Replace feed-forward MLP with Gated Transformer-XL for temporal context — "
     "handle wind disturbances and improve trajectory planning during descent. "
     "Recent aerospace results (Federici et al., Carradori et al.) show promise.",
     MEDIUM_BLUE),
    ("Domain Randomization",
     "Randomize aerodynamic coefficients, servo parameters, mass properties, "
     "and sensor noise during training for robust sim-to-real transfer.",
     GREEN),
    ("Lateral Guidance",
     "Potential-based reward shaping (Ng et al., 1999), hierarchical policies "
     "(separate lateral/vertical controllers), extended curriculum strategies.",
     ACCENT_ORANGE),
    ("Hardware Validation",
     "Deploy trained policies on the physical EDF drone testbed. "
     "Validate sim-to-real transfer and characterize the reality gap.",
     RED),
    ("Wind Disturbance Rejection",
     "Train and evaluate under stochastic wind conditions using the simulation's "
     "wind model with 27-dim observation space including wind estimates.",
     RGBColor(0x7B, 0x1F, 0xA2)),
]

for i, (title, desc, color) in enumerate(future_items):
    y = Inches(1.2) + Inches(i * 1.15)
    # Color bar
    add_rect(slide, Inches(0.5), y, Inches(0.12), Inches(0.9), color)
    add_textbox(slide, Inches(0.9), y, Inches(4.0), Inches(0.35),
                title, font_size=18, color=color, bold=True)
    add_textbox(slide, Inches(0.9), y + Inches(0.35), Inches(11.5), Inches(0.55),
                desc, font_size=14, color=DARK_GRAY, line_spacing=1.2)

# Timeline hint
add_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
            "Near-term priority: GTrXL integration + domain randomization → hardware flight tests. "
            "The 24-dim observation space is already compatible with transformer sequence input.",
            font_size=15, color=DARK_BLUE, line_spacing=1.2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion & Thank You
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Top blue band
add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(2.5), DARK_BLUE)
add_rect(slide, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.08), ACCENT_ORANGE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.6),
            "Conclusion", font_size=34, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.2),
            "PPO can learn thrust-vectored fin control for an EDF VTOL vehicle in simulation — "
            "achieving 100% landing rate with soft touchdowns, opening the path to sim-to-real transfer.",
            font_size=20, color=RGBColor(0xBB, 0xD5, 0xED), line_spacing=1.3)

# Key contributions
add_section_header(slide, Inches(0.6), Inches(2.8), Inches(12.0),
                   "KEY CONTRIBUTIONS")

contributions = [
    "GPU-accelerated Isaac Sim environment with semi-empirical aerodynamic, propulsion, and actuator models",
    "PPO hover policy: 0.085 m mean error in only 2,048 steps (residual-PID architecture)",
    "PPO landing policy: 100% success rate, 0.126 m/s mean touchdown speed over 5M steps",
    "Critical insight: initialization priors on action distribution prevent degenerate local optima",
    "Reward design principle: terminal magnitudes must dominate integrated per-step costs",
]

for i, contrib in enumerate(contributions):
    y = Inches(3.2) + Inches(i * 0.55)
    # Checkmark
    add_textbox(slide, Inches(0.6), y, Inches(0.4), Inches(0.4),
                "✓", font_size=18, color=GREEN, bold=True)
    add_textbox(slide, Inches(1.0), y, Inches(11.5), Inches(0.45),
                contrib, font_size=16, color=DARK_GRAY, line_spacing=1.15)

# Contact / repo
add_rect(slide, Inches(0), Inches(6.2), SLIDE_WIDTH, Inches(1.3), LIGHT_GRAY)
add_textbox(slide, Inches(0.8), Inches(6.3), Inches(6), Inches(0.4),
            "Tang Zijian (Jacob Tang)  |  Minnesota State University, Mankato",
            font_size=14, color=DARK_GRAY)
add_textbox(slide, Inches(0.8), Inches(6.7), Inches(10), Inches(0.4),
            "Code & Data: github.com/Jacob19999/Transformer-rl-retro-propulsion",
            font_size=14, color=MEDIUM_BLUE, bold=True)

add_textbox(slide, Inches(9.0), Inches(6.3), Inches(4.0), Inches(0.8),
            "Thank You!\nQuestions?",
            font_size=28, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.RIGHT)

add_footer_bar(slide)
add_slide_number(slide, 20)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "TVC_RL_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
