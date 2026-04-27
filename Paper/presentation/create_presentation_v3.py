"""
TVC RL Presentation — v3
~20 minutes, expanded scripts, 22 slides.
All math rendered via matplotlib mathtext → PNG → embedded in PPTX.
Speaker notes embedded in every slide's notes pane.
"""

import os, io
import matplotlib
matplotlib.use("Agg")
matplotlib.rcParams.update({
    "text.usetex": False,
    "mathtext.fontset": "cm",
    "font.family": "serif",
    "figure.facecolor": "none",
    "axes.facecolor": "none",
})
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ────────────────────────────────────────────────────
DARK_BLUE    = RGBColor(0x00, 0x3B, 0x6F)
MEDIUM_BLUE  = RGBColor(0x00, 0x5A, 0x9C)
LIGHT_BLUE   = RGBColor(0xD6, 0xEA, 0xF8)
ACCENT       = RGBColor(0xE8, 0x6C, 0x00)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MED_GRAY     = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
GREEN        = RGBColor(0x00, 0x7A, 0x33)
RED          = RGBColor(0xC0, 0x39, 0x2B)
WARM_YEL     = RGBColor(0xFF, 0xF3, 0xE0)
WARM_RED     = RGBColor(0xFC, 0xE4, 0xEC)
WARM_GRN     = RGBColor(0xE8, 0xF5, 0xE9)

SW = Inches(13.333)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

HERE     = os.path.dirname(__file__)
FIG_DIR  = os.path.join(os.path.dirname(HERE), "figures")
TOTAL    = 22   # total slide count

# ── Math renderer ───────────────────────────────────────────────
def render_eq(s, fs=22, color="#003B6F", bg=None, dpi=200):
    fig = plt.figure(figsize=(0.01, 0.01))
    ax  = fig.add_axes([0, 0, 1, 1]); ax.set_axis_off()
    if bg: fig.patch.set_facecolor(bg); ax.set_facecolor(bg)
    else:  fig.patch.set_alpha(0);     ax.patch.set_alpha(0)
    t = ax.text(0.5, 0.5, s, fontsize=fs, color=color,
                ha="center", va="center", transform=ax.transAxes)
    fig.canvas.draw()
    rend = fig.canvas.get_renderer()
    bb   = t.get_window_extent(renderer=rend)
    w_in = (bb.width  / dpi) + 0.3
    h_in = (bb.height / dpi) + 0.24
    fig.set_size_inches(max(w_in, 0.5), max(h_in, 0.3))
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi,
                bbox_inches="tight", transparent=(bg is None), pad_inches=0.05)
    plt.close(fig); buf.seek(0)
    return buf

def eq(sl, s, l, t, h=Inches(0.55), fs=22, color="#003B6F", bg=None):
    buf = render_eq(s, fs=fs, color=color, bg=bg)
    return sl.shapes.add_picture(buf, l, t, height=h)

# ── Layout helpers ──────────────────────────────────────────────
def bg(sl, c): f=sl.background.fill; f.solid(); f.fore_color.rgb=c

def rect(sl, l, t, w, h, fc, lc=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=fc
    if lc: s.line.color.rgb=lc
    else:  s.line.fill.background()
    return s

def tb(sl, l, t, w, h, text, fs=18, c=DARK_GRAY, bold=False,
       align=PP_ALIGN.LEFT, font="Calibri", ls=1.2):
    box=sl.shapes.add_textbox(l,t,w,h)
    tf=box.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text
    p.font.size=Pt(fs); p.font.color.rgb=c
    p.font.bold=bold; p.font.name=font
    p.alignment=align
    p.space_after=Pt(0); p.space_before=Pt(0)
    if ls!=1.0: p.line_spacing=Pt(fs*ls)
    return box

def bul(sl, l, t, w, h, items, fs=20, c=DARK_GRAY,
        font="Calibri", sp=8, ls=1.15):
    box=sl.shapes.add_textbox(l,t,w,h)
    tf=box.text_frame; tf.word_wrap=True
    for i,item in enumerate(items):
        text,lv=(item if isinstance(item,tuple) else (item,0))
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        bch="●" if lv==0 else "○"
        p.text="    "*lv+bch+"  "+text
        p.font.size=Pt(fs-lv*2); p.font.color.rgb=c
        p.font.name=font; p.space_after=Pt(sp); p.space_before=Pt(0)
        if ls!=1.0: p.line_spacing=Pt((fs-lv*2)*ls)
    return box

def lbl(sl, l, t, w, text, fs=14, c=MEDIUM_BLUE):
    return tb(sl,l,t,w,Inches(0.4),text,fs=fs,c=c,bold=True)

def foot(sl):
    rect(sl,Inches(0),Inches(7.25),SW,Inches(0.25),DARK_BLUE)

def snum(sl, n):
    tb(sl,Inches(12.2),Inches(7.05),Inches(1.0),Inches(0.35),
       f"{n}/{TOTAL}",fs=11,c=MED_GRAY,align=PP_ALIGN.RIGHT)

def hdr(sl, title, n):
    rect(sl,Inches(0),Inches(0),SW,Inches(0.95),DARK_BLUE)
    tb(sl,Inches(0.6),Inches(0.15),Inches(10),Inches(0.7),
       title,fs=30,c=WHITE,bold=True)
    foot(sl); snum(sl,n)

def oval(sl, x, y, txt, bg_c=MEDIUM_BLUE):
    s=sl.shapes.add_shape(MSO_SHAPE.OVAL,x,y,Inches(0.55),Inches(0.55))
    s.fill.solid(); s.fill.fore_color.rgb=bg_c; s.line.fill.background()
    tf=s.text_frame
    tf.paragraphs[0].text=txt
    tf.paragraphs[0].font.size=Pt(18)
    tf.paragraphs[0].font.color.rgb=WHITE
    tf.paragraphs[0].font.bold=True
    tf.paragraphs[0].alignment=PP_ALIGN.CENTER
    tf.vertical_anchor=MSO_ANCHOR.MIDDLE

def img(sl, path, l, t, w=None, h=None):
    if os.path.exists(path):
        kw={}
        if w: kw["width"]=w
        if h: kw["height"]=h
        sl.shapes.add_picture(path,l,t,**kw)

def mktbl(sl, l, t, w, h, rows, cols, data,
          cws=None, hc=DARK_BLUE, fs=14):
    ts=sl.shapes.add_table(rows,cols,l,t,w,h)
    tbl=ts.table
    if cws:
        for i,cw in enumerate(cws): tbl.columns[i].width=cw
    for r in range(rows):
        for c in range(cols):
            cell=tbl.cell(r,c); cell.text=str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(fs); p.font.name="Calibri"
                if r==0:
                    p.font.bold=True; p.font.color.rgb=WHITE
                    p.alignment=PP_ALIGN.CENTER
                else:
                    p.font.color.rgb=DARK_GRAY
                    p.alignment=PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER
            if r==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=hc
            elif r%2==0:
                cell.fill.solid(); cell.fill.fore_color.rgb=LIGHT_GRAY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
    return ts

def notes(sl, text):
    sl.notes_slide.notes_text_frame.text = text

