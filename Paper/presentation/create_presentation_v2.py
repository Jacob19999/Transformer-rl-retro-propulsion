"""
TVC RL Presentation — v2
All math expressions rendered via matplotlib mathtext → PNG → embedded in PPTX.
"""

import os
import io
import matplotlib
matplotlib.use("Agg")
matplotlib.rcParams.update({
    "text.usetex": False,
    "mathtext.fontset": "cm",       # Computer Modern — closest to LaTeX look
    "font.family": "serif",
    "figure.facecolor": "none",
    "axes.facecolor": "none",
})
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────
DARK_BLUE    = RGBColor(0x00, 0x3B, 0x6F)
MEDIUM_BLUE  = RGBColor(0x00, 0x5A, 0x9C)
LIGHT_BLUE   = RGBColor(0xD6, 0xEA, 0xF8)
ACCENT_ORANGE= RGBColor(0xE8, 0x6C, 0x00)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY  = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
GREEN        = RGBColor(0x00, 0x7A, 0x33)
RED          = RGBColor(0xC0, 0x39, 0x2B)
WARM_YELLOW  = RGBColor(0xFF, 0xF3, 0xE0)
WARM_RED     = RGBColor(0xFC, 0xE4, 0xEC)
WARM_GREEN   = RGBColor(0xE8, 0xF5, 0xE9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

ASSETS_DIR = os.path.join(os.path.dirname(__file__), "_eq_assets")
os.makedirs(ASSETS_DIR, exist_ok=True)

FIGURES_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "figures")

# ── Math image renderer ─────────────────────────────────────────

def render_eq(latex_str, fontsize=22, text_color="#003B6F", bg_color=None,
              pad_w=0.15, pad_h=0.12, dpi=200):
    """
    Render a LaTeX math string (mathtext) to a PNG byte-stream.
    Returns a BytesIO object ready for pptx add_picture.
    """
    fig = plt.figure(figsize=(0.01, 0.01))  # will be resized by tight_layout
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_axis_off()
    if bg_color:
        fig.patch.set_facecolor(bg_color)
        ax.set_facecolor(bg_color)
    else:
        fig.patch.set_alpha(0)
        ax.patch.set_alpha(0)

    t = ax.text(0.5, 0.5, latex_str,
                fontsize=fontsize,
                color=text_color,
                ha="center", va="center",
                transform=ax.transAxes)

    # Measure bounding box and resize figure
    fig.canvas.draw()
    renderer = fig.canvas.get_renderer()
    bb = t.get_window_extent(renderer=renderer)
    w_in = (bb.width  / dpi) + 2 * pad_w
    h_in = (bb.height / dpi) + 2 * pad_h
    fig.set_size_inches(max(w_in, 0.5), max(h_in, 0.3))

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi,
                bbox_inches="tight", transparent=(bg_color is None),
                pad_inches=0.05)
    plt.close(fig)
    buf.seek(0)
    return buf


def add_eq(slide, latex_str, left, top, height=Inches(0.55),
           fontsize=22, text_color="#003B6F", bg_color=None):
    """Render equation and insert into slide, auto-sizing width."""
    buf = render_eq(latex_str, fontsize=fontsize,
                    text_color=text_color, bg_color=bg_color)
    pic = slide.shapes.add_picture(buf, left, top, height=height)
    return pic


# ── Generic helpers ─────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color: s.line.color.rgb = line_color
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, fs=18, color=DARK_GRAY, bold=False,
       align=PP_ALIGN.LEFT, font="Calibri", ls=1.2):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(fs); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font
    p.alignment = align
    p.space_after = Pt(0); p.space_before = Pt(0)
    if ls != 1.0: p.line_spacing = Pt(fs * ls)
    return box

def bullets(slide, l, t, w, h, items, fs=20, color=DARK_GRAY,
            font="Calibri", spacing=8, ls=1.15):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = "●" if level == 0 else "○"
        p.text = "    " * level + bullet + "  " + text
        p.font.size = Pt(fs - level * 2)
        p.font.color.rgb = color; p.font.name = font
        p.space_after = Pt(spacing); p.space_before = Pt(0)
        if ls != 1.0: p.line_spacing = Pt((fs - level * 2) * ls)
    return box

def section_label(slide, l, t, w, text, fs=14, color=MEDIUM_BLUE):
    return tb(slide, l, t, w, Inches(0.4), text, fs=fs, color=color, bold=True)

def footer(slide):
    add_rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.25), DARK_BLUE)

def slide_num(slide, n, total=20):
    tb(slide, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.35),
       f"{n}/{total}", fs=11, color=MEDIUM_GRAY, align=PP_ALIGN.RIGHT)

def header(slide, title, n, total=20):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.95), DARK_BLUE)
    tb(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
       title, fs=30, color=WHITE, bold=True)
    footer(slide); slide_num(slide, n, total)

def make_table(slide, l, t, w, h, rows, cols, data,
               col_widths=None, hdr_color=DARK_BLUE, fs=14):
    ts = slide.shapes.add_table(rows, cols, l, t, w, h)
    tbl = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths): tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c); cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(fs); p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True; p.font.color.rgb = WHITE
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = DARK_GRAY
                    p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = hdr_color
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    return ts

def oval_num(slide, x, y, num_str, bg=MEDIUM_BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.55), Inches(0.55))
    s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
    tf = s.text_frame
    tf.paragraphs[0].text = num_str
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def img(slide, path, l, t, w=None, h=None):
    """Insert an image file if it exists."""
    if os.path.exists(path):
        if w and h: slide.shapes.add_picture(path, l, t, width=w, height=h)
        elif h:     slide.shapes.add_picture(path, l, t, height=h)
        elif w:     slide.shapes.add_picture(path, l, t, width=w)
        else:       slide.shapes.add_picture(path, l, t)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(3.2), DARK_BLUE)
add_rect(sl, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08), ACCENT_ORANGE)

tb(sl, Inches(0.8), Inches(0.55), Inches(11.7), Inches(1.3),
   "Reinforcement Learning for Thrust-Vectored\nFin Control of an EDF VTOL Vehicle",
   fs=36, color=WHITE, bold=True, ls=1.15)
tb(sl, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.8),
   "PPO-Based Attitude and Position Control in GPU-Accelerated Simulation",
   fs=22, color=RGBColor(0xBB, 0xD5, 0xED))
tb(sl, Inches(0.8), Inches(3.8), Inches(6), Inches(0.5),
   "Tang Zijian (Jacob Tang)", fs=22, color=DARK_BLUE, bold=True)
tb(sl, Inches(0.8), Inches(4.35), Inches(6), Inches(0.8),
   "Department of Applied Data Science\nMinnesota State University, Mankato",
   fs=16, color=MEDIUM_GRAY)
tb(sl, Inches(0.8), Inches(5.5), Inches(6), Inches(0.4),
   "AI Conference 2026", fs=16, color=MEDIUM_GRAY, bold=True)

# Drone image (right side)
cad_path = os.path.join(FIGURES_DIR, "cad_v1.png")
if os.path.exists(cad_path):
    img(sl, cad_path, Inches(8.0), Inches(3.5), h=Inches(3.3))
else:
    add_rect(sl, Inches(8.0), Inches(3.6), Inches(4.5), Inches(3.2), LIGHT_GRAY)
    tb(sl, Inches(8.2), Inches(4.6), Inches(4.1), Inches(1.0),
       "[EDF Drone CAD Render]", fs=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

footer(sl); slide_num(sl, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Outline
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Presentation Outline", 2)

agenda = [
    ("1", "Motivation & Problem Statement",    "Why thrust-vectored EDF drones?"),
    ("2", "Vehicle Platform",                  "Physical testbed specifications"),
    ("3", "Simulation Environment",            "GPU-accelerated Isaac Sim setup"),
    ("4", "Physics Models",                    "Propulsion, aerodynamics, actuators"),
    ("5", "PPO Algorithm & Training",          "Observations, actions, reward design"),
    ("6", "Experimental Results",              "Hover & landing performance"),
    ("7", "Discussion & Future Work",          "Challenges and next steps"),
]
for i, (num, title, desc) in enumerate(agenda):
    y = Inches(1.3) + Inches(i * 0.78)
    oval_num(sl, Inches(1.0), y, num)
    tb(sl, Inches(1.8), y - Inches(0.02), Inches(5), Inches(0.35),
       title, fs=20, color=DARK_BLUE, bold=True)
    tb(sl, Inches(1.8), y + Inches(0.30), Inches(5), Inches(0.30),
       desc, fs=14, color=MEDIUM_GRAY)
    if i < len(agenda) - 1:
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1.25), y + Inches(0.55),
                                   Inches(0.05), Inches(0.23))
        bar.fill.solid(); bar.fill.fore_color.rgb = LIGHT_BLUE
        bar.line.fill.background()

tb(sl, Inches(8.5), Inches(1.5), Inches(4), Inches(0.4),
   "~20 minute presentation", fs=16, color=ACCENT_ORANGE, bold=True)
add_rect(sl, Inches(8.0), Inches(2.2), Inches(4.8), Inches(3.5), LIGHT_BLUE)
tb(sl, Inches(8.3), Inches(2.4), Inches(4.2), Inches(0.4),
   "Key Takeaway", fs=18, color=DARK_BLUE, bold=True)
tb(sl, Inches(8.3), Inches(2.9), Inches(4.2), Inches(2.5),
   "Can deep RL learn to control a vehicle that steers by deflecting fins "
   "in its own exhaust stream?\n\nYes — PPO achieves 100% landing rate with "
   "soft touchdowns in simulation, opening the door to sim-to-real transfer.",
   fs=15, color=DARK_GRAY, ls=1.3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Motivation
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Motivation: Why Thrust-Vectored EDF Drones?", 3)

section_label(sl, Inches(0.6), Inches(1.15), Inches(5.5), "THE CONTROL CHALLENGE")
bullets(sl, Inches(0.6), Inches(1.55), Inches(5.8), Inches(4.5), [
    "Conventional multirotors: independent rotor speeds → decoupled control",
    "Thrust-vectored EDF: single thrust source + 4 jet-vane fins",
    "Fin effectiveness ∝ throttle² — at 50% throttle, only 25% authority",
    "Nonlinear aerodynamics, actuator lag, gyroscopic precession",
    "Classical PID/LQR limited to narrow operating envelope",
    "RL learns the full nonlinear mapping end-to-end",
], fs=18, spacing=6)

add_rect(sl, Inches(7.0), Inches(1.15), Inches(5.8), Inches(2.8), LIGHT_BLUE)
tb(sl, Inches(7.3), Inches(1.3), Inches(5.2), Inches(0.4),
   "Why It Matters", fs=18, color=DARK_BLUE, bold=True)
bullets(sl, Inches(7.3), Inches(1.75), Inches(5.2), Inches(2.0), [
    "Aerodynamic efficiency over multirotors",
    "Noise reduction (single fan vs. 4+ rotors)",
    "Mechanical simplicity for certain missions",
    "Applicable to rocket landing (jet vanes)",
], fs=16, spacing=6)

# Equation box — rendered math
add_rect(sl, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.5), WARM_YELLOW)
tb(sl, Inches(7.3), Inches(4.35), Inches(5.2), Inches(0.35),
   "The Core Coupling", fs=16, color=ACCENT_ORANGE, bold=True)
add_eq(sl, r"$F_{\rm fin} = q \cdot S \cdot C_N(\alpha)$",
       Inches(7.4), Inches(4.8), height=Inches(0.55), fontsize=24)
add_eq(sl, r"$q \;\propto\; u_{\rm thr}^{2}$",
       Inches(7.4), Inches(5.45), height=Inches(0.5), fontsize=22)
tb(sl, Inches(7.3), Inches(6.05), Inches(5.2), Inches(0.5),
   "→ Reducing throttle to descend simultaneously reduces lateral steering authority",
   fs=14, color=DARK_GRAY, ls=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Related Work
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Related Work & Research Gap", 4)

cols_data = [
    ("RL for Aerial Vehicles", [
        "Hwangbo et al. — PPO for agile quadrotor control",
        "Koch et al. — Deep RL for attitude control",
        "Pi et al. — Sim-to-real under wind disturbances",
        "GPU-parallel training (Isaac Gym, Rudin et al.)",
    ]),
    ("Thrust Vector Control", [
        "Long history in rocketry (gimbaled nozzles)",
        "Jet-vane TVC for VTOL vehicles",
        "Classical: PID/LQR around hover linearization",
        "Limited envelope for nonlinear regime",
    ]),
    ("Sim-to-Real Transfer", [
        "Domain randomization (Tobin et al.)",
        "System identification (Yu et al.)",
        "CAD2RL: simulation-only to real flight",
        "Fidelity of dynamics model is critical",
    ]),
]
for i, (title, items) in enumerate(cols_data):
    x = Inches(0.5) + Inches(i * 4.2)
    add_rect(sl, x, Inches(1.15), Inches(3.9), Inches(0.5), MEDIUM_BLUE)
    tb(sl, x + Inches(0.15), Inches(1.2), Inches(3.6), Inches(0.4),
       title, fs=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    bullets(sl, x + Inches(0.15), Inches(1.8), Inches(3.6), Inches(3.5),
            items, fs=14, spacing=6, ls=1.2)

add_rect(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.3), WARM_YELLOW)
tb(sl, Inches(0.8), Inches(5.6), Inches(2.0), Inches(0.4),
   "RESEARCH GAP", fs=16, color=ACCENT_ORANGE, bold=True)
tb(sl, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6),
   "RL for thrust-vectored fin control of EDF vehicles remains largely unexplored. "
   "Key distinctions: single thrust source with coupled fins, quadratic throttle-authority "
   "dependence, hobby-grade actuator dynamics, and gyroscopic effects from high-speed rotor.",
   fs=15, color=DARK_GRAY, ls=1.25)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Vehicle Platform
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Vehicle Platform: EDF TVC Drone Testbed", 5)

section_label(sl, Inches(0.6), Inches(1.15), Inches(5.5), "PHYSICAL SPECIFICATIONS")
spec_data = [
    ["Parameter", "Value"],
    ["Total mass", "3.1 kg"],
    ["Body length / diameter", "0.35 m / 0.12 m"],
    ["Ixx, Iyy  /  Izz", "0.05 / 0.02 kg·m²"],
    ["EDF unit", "FMS 90 mm 12-blade"],
    ["Motor", "4068-KV1850 inrunner"],
    ["Battery", "6S LiPo"],
    ["Servos", "4× MG996R"],
    ["Fin count / area", "4 / 0.002 m² each"],
    ["Max fin deflection", "±15° (0.262 rad)"],
    ["Max thrust", "39.2 N"],
]
make_table(sl, Inches(0.6), Inches(1.55), Inches(5.5), Inches(4.8),
           len(spec_data), 2, spec_data,
           col_widths=[Inches(3.2), Inches(2.3)], fs=13)

cad_path = os.path.join(FIGURES_DIR, "cad_v1.png")
if os.path.exists(cad_path):
    img(sl, cad_path, Inches(7.0), Inches(1.15), h=Inches(3.0))
else:
    add_rect(sl, Inches(7.0), Inches(1.15), Inches(5.8), Inches(2.8), LIGHT_GRAY)
    tb(sl, Inches(7.5), Inches(2.2), Inches(4.8), Inches(0.5),
       "[CAD Render: cad_v1.png]", fs=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(7.0), Inches(4.3), Inches(5.8), Inches(2.5), LIGHT_BLUE)
tb(sl, Inches(7.3), Inches(4.45), Inches(5.2), Inches(0.35),
   "Control Architecture", fs=16, color=DARK_BLUE, bold=True)
bullets(sl, Inches(7.3), Inches(4.85), Inches(5.2), Inches(1.8), [
    "Forward fin (+X): pitch control",
    "Right fin (+Y): roll control",
    "Aft fin (−X): complementary pitch",
    "Left fin (−Y): complementary roll",
    "Yaw: differential tangential forces",
    "Altitude: EDF throttle",
    "→ 5 actuators → full 6-DOF control",
], fs=14, spacing=4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Simulation Environment
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Simulation Environment: NVIDIA Isaac Sim", 6)

bullets(sl, Inches(0.6), Inches(1.25), Inches(5.8), Inches(5.5), [
    "Isaac Lab v2.3.2 — PhysX 5 GPU-accelerated rigid-body physics",
    "128 parallel environments on a single GPU",
    "Gymnasium-compatible vectorized interface",
    "Physics timestep: 8.33 ms (120 Hz)",
    "RL decision timestep: 33.3 ms (30 Hz, decimation = 4)",
    "Semi-empirical dynamics calibrated to hardware",
    "Two force dispatch modes:",
    ("Per-link: forces at center-of-pressure locations", 1),
    ("Collapsed body wrench: for debugging/validation", 1),
], fs=18, spacing=6)

env_path = os.path.join(FIGURES_DIR, "isaacsim_128_env.png")
if os.path.exists(env_path):
    img(sl, env_path, Inches(7.0), Inches(1.15), h=Inches(3.5))
else:
    add_rect(sl, Inches(7.0), Inches(1.15), Inches(5.8), Inches(3.5), LIGHT_GRAY)
    tb(sl, Inches(7.5), Inches(2.5), Inches(4.8), Inches(0.5),
       "[isaacsim_128_env.png]", fs=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(7.0), Inches(4.9), Inches(5.8), Inches(1.8), LIGHT_BLUE)
tb(sl, Inches(7.3), Inches(5.0), Inches(5.2), Inches(0.35),
   "Key Numbers", fs=16, color=DARK_BLUE, bold=True)
for i, (val, label) in enumerate([("128", "Parallel Envs"), ("120 Hz", "Physics"),
                                    ("30 Hz", "RL Rate"), ("4×", "Decimation")]):
    x = Inches(7.3) + Inches(i * 1.3)
    tb(sl, x, Inches(5.45), Inches(1.2), Inches(0.5),
       val, fs=22, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x, Inches(5.95), Inches(1.2), Inches(0.4),
       label, fs=11, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Propulsion Model
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Propulsion Model: EDF Dynamics", 7)

section_label(sl, Inches(0.6), Inches(1.15), Inches(6.0),
              "FIRST-ORDER SPOOL + QUADRATIC THRUST")

# Rendered equations
add_eq(sl, r"$\dot{\omega} = \dfrac{\omega_{\rm target} - \omega}{\tau_{\rm motor}}$",
       Inches(0.8), Inches(1.6), height=Inches(0.7), fontsize=26)
tb(sl, Inches(5.5), Inches(1.75), Inches(2.5), Inches(0.4),
   "← spool dynamics", fs=14, color=MEDIUM_GRAY)

add_eq(sl, r"$T = k_T \,\omega^{2}$",
       Inches(0.8), Inches(2.5), height=Inches(0.6), fontsize=26)
tb(sl, Inches(5.5), Inches(2.65), Inches(2.5), Inches(0.4),
   "← quadratic thrust", fs=14, color=MEDIUM_GRAY)

section_label(sl, Inches(0.6), Inches(3.3), Inches(6.0), "THREE TORQUE COMPONENTS")
bullets(sl, Inches(0.6), Inches(3.7), Inches(6.0), Inches(3.0), [
    "Static reaction torque: opposes rotor spin",
    "Dynamic spool torque: reaction to rotor acceleration",
    "Gyroscopic precession: couples body rotation to rotor angular momentum",
    ("Critical for realistic yaw dynamics", 1),
], fs=16, spacing=5)

# Gyro equation
add_eq(sl,
       r"$\boldsymbol{\tau}_{\rm gyro} = -\boldsymbol{\omega}_{\rm body} \times (I_{\rm rotor}\,\omega\,\hat{e}_z)$",
       Inches(0.8), Inches(5.5), height=Inches(0.6), fontsize=22)

# Right: parameters table
section_label(sl, Inches(7.0), Inches(1.15), Inches(5.5), "EDF PARAMETERS")
edf_data = [
    ["Parameter", "Value"],
    ["Max thrust T_max", "39.2 N"],
    ["Motor time constant τ", "0.15 s"],
    ["Max angular velocity ω_max", "4300 rad/s"],
    ["Rotor inertia I_rotor", "0.0002 kg·m²"],
    ["Exhaust speed (nominal)", "116 m/s"],
    ["Gyro torque scale", "0.1"],
]
make_table(sl, Inches(7.0), Inches(1.55), Inches(5.5), Inches(3.0),
           len(edf_data), 2, edf_data,
           col_widths=[Inches(3.2), Inches(2.3)], fs=13)

add_rect(sl, Inches(7.0), Inches(4.9), Inches(5.8), Inches(1.8), WARM_YELLOW)
tb(sl, Inches(7.3), Inches(5.0), Inches(5.2), Inches(0.35),
   "Key Insight", fs=16, color=ACCENT_ORANGE, bold=True)
tb(sl, Inches(7.3), Inches(5.4), Inches(5.2), Inches(1.0),
   "Gyroscopic precession couples body rotation to rotor angular momentum — "
   "a real physical effect that makes yaw control non-trivial and distinguishes "
   "this from simplified point-mass models.",
   fs=14, color=DARK_GRAY, ls=1.3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — Fin Aerodynamic Model
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Fin Aerodynamic Model: Jet-Vane Forces", 8)

section_label(sl, Inches(0.6), Inches(1.15), Inches(6.0), "SEMI-EMPIRICAL FORCE MODEL")

eq_specs = [
    (r"$q = \frac{1}{2}\rho\,v_{\rm exh}^{2}\cdot u_{\rm thr}^{2}\cdot k_{\rm duct}$",
     "dynamic pressure"),
    (r"$C_N(\alpha) = C_{N_\alpha}\cdot\alpha\cdot(1 - k_{\rm sat}\,\alpha^2)$",
     "normal force coeff"),
    (r"$C_D(\alpha) = C_{D_0} + C_{D_{\alpha^2}}\,\alpha^2$",
     "drag coeff"),
    (r"$F_n = q\,S_{\rm fin}\,C_N(\alpha),\quad F_t = q\,S_{\rm fin}\,C_D(\alpha)$",
     "per-fin forces"),
]
for i, (eq_str, label) in enumerate(eq_specs):
    y = Inches(1.6) + Inches(i * 1.1)
    tb(sl, Inches(0.8), y, Inches(2.5), Inches(0.35),
       label + ":", fs=13, color=MEDIUM_BLUE, bold=True)
    add_eq(sl, eq_str, Inches(0.8), y + Inches(0.3), height=Inches(0.6), fontsize=20)

# Right: parameters
aero_data = [
    ["Parameter", "Value"],
    ["Air density ρ", "1.225 kg/m³"],
    ["Exhaust speed v_exh", "116 m/s"],
    ["Duct correction k_duct", "1.3"],
    ["Lift slope C_Nα", "3.5 rad⁻¹"],
    ["Stall saturation k_sat", "2.0"],
    ["Zero-defl drag C_D0", "0.05"],
    ["Drag slope C_Dα²", "1.5"],
    ["Fin area S_fin", "0.002 m²"],
]
make_table(sl, Inches(7.0), Inches(1.15), Inches(5.5), Inches(4.0),
           len(aero_data), 2, aero_data,
           col_widths=[Inches(3.2), Inches(2.3)], fs=13)

add_rect(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.3), WARM_RED)
tb(sl, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.35),
   "CRITICAL COUPLING", fs=16, color=RED, bold=True)
tb(sl, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.7),
   "Fin force ∝ throttle² — at 50% throttle, fin forces drop to 25% of maximum. "
   "This quadratic coupling is THE defining challenge of jet-vane TVC and why "
   "classical linear controllers struggle outside the hover envelope.",
   fs=15, color=DARK_GRAY, ls=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — PPO Algorithm
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Learning Algorithm: Proximal Policy Optimization", 9)

section_label(sl, Inches(0.6), Inches(1.15), Inches(6.0),
              "PPO CLIPPED SURROGATE OBJECTIVE")

add_eq(sl,
       r"$L^{\rm CLIP}(\theta) = \hat{\mathbb{E}}_t\!\left[\min\!\left(r_t(\theta)\hat{A}_t,\;"
       r"\mathrm{clip}(r_t(\theta),1{-}\epsilon,1{+}\epsilon)\hat{A}_t\right)\right]$",
       Inches(0.7), Inches(1.6), height=Inches(0.7), fontsize=20)

section_label(sl, Inches(0.6), Inches(2.5), Inches(6.0), "GENERALIZED ADVANTAGE ESTIMATE")
add_eq(sl,
       r"$\hat{A}_t = \sum_{l=0}^{T-t}(\gamma\lambda)^l\,\delta_{t+l},"
       r"\quad \delta_t = r_t + \gamma V(s_{t+1}) - V(s_t)$",
       Inches(0.7), Inches(2.9), height=Inches(0.65), fontsize=20)

section_label(sl, Inches(0.6), Inches(3.7), Inches(6.0), "COMBINED LOSS")
add_eq(sl,
       r"$L(\theta) = L^{\rm CLIP}(\theta) - c_1 L^{\rm VF}(\theta) + c_2 H[\pi_\theta]$",
       Inches(0.7), Inches(4.1), height=Inches(0.6), fontsize=20)

bullets(sl, Inches(0.6), Inches(4.9), Inches(6.0), Inches(1.8), [
    "On-policy actor-critic — trust-region constraint prevents destructive updates",
    "Clip ε = 0.2,  γ = 0.99,  λ = 0.95,  KL_target = 0.03",
    "Value loss c₁ = 0.5,  Entropy bonus c₂ = 0.005",
], fs=16, spacing=6)

# Right: hyperparameters
section_label(sl, Inches(7.0), Inches(1.15), Inches(5.5), "HYPERPARAMETERS (LANDING TASK)")
hyper_data = [
    ["Parameter", "Value"],
    ["Parallel environments", "128"],
    ["Rollout length", "128 steps"],
    ["Minibatches / update", "8"],
    ["Update epochs", "4"],
    ["Learning rate", "3 × 10⁻⁴"],
    ["Discount γ", "0.99"],
    ["GAE λ", "0.95"],
    ["Clip ε", "0.2"],
    ["Entropy coeff", "0.005"],
    ["Max grad norm", "0.5"],
    ["Target KL", "0.03"],
]
make_table(sl, Inches(7.0), Inches(1.55), Inches(5.5), Inches(5.2),
           len(hyper_data), 2, hyper_data,
           col_widths=[Inches(3.2), Inches(2.3)], fs=13)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — Observation & Action Space
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Observation & Action Space Design", 10)

section_label(sl, Inches(0.6), Inches(1.15), Inches(5.5),
              "24-DIMENSIONAL OBSERVATION VECTOR")
obs_data = [
    ["Index", "Component", "Dim"],
    ["0–2",   "Position error (m)",              "3"],
    ["3–6",   "Attitude quaternion [w,x,y,z]",   "4"],
    ["7–9",   "Linear velocity, body FRD",        "3"],
    ["10–12", "Angular velocity, body FRD",       "3"],
    ["13",    "Height above ground",              "1"],
    ["14–17", "Fin angles (rad)",                 "4"],
    ["18–21", "Fin angular rates (rad/s)",        "4"],
    ["22",    "Motor RPM, normalized",            "1"],
    ["23",    "Contact state",                    "1"],
]
make_table(sl, Inches(0.5), Inches(1.55), Inches(6.2), Inches(4.2),
           len(obs_data), 3, obs_data,
           col_widths=[Inches(1.0), Inches(3.8), Inches(1.4)], fs=13)

section_label(sl, Inches(7.0), Inches(1.15), Inches(5.5),
              "5-DIMENSIONAL ACTION VECTOR")
add_eq(sl,
       r"$\mathbf{a} = [\alpha_1,\,\alpha_2,\,\alpha_3,\,\alpha_4,\,u_{\rm thr}]$",
       Inches(7.1), Inches(1.6), height=Inches(0.6), fontsize=22)

bullets(sl, Inches(7.0), Inches(2.4), Inches(5.5), Inches(2.0), [
    "4 fin deflections: ±15° (±0.262 rad)",
    "1 throttle command: [0, 1]",
    "Network outputs tanh → linear scaling",
    "Smaller action space than quadrotor (4 rotors)",
], fs=16, spacing=6)

add_rect(sl, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.6), LIGHT_BLUE)
tb(sl, Inches(7.3), Inches(4.35), Inches(5.2), Inches(0.35),
   "Network Architecture", fs=16, color=DARK_BLUE, bold=True)
tb(sl, Inches(7.3), Inches(4.8), Inches(5.2), Inches(0.4),
   "Actor:  Obs(24) → FC(256) → tanh → FC(256) → tanh → FC(5)",
   fs=13, color=DARK_GRAY, font="Consolas")
tb(sl, Inches(7.3), Inches(5.25), Inches(5.2), Inches(0.4),
   "Critic: Obs(24) → FC(256) → tanh → FC(256) → tanh → FC(1)",
   fs=13, color=DARK_GRAY, font="Consolas")
tb(sl, Inches(7.3), Inches(5.75), Inches(5.2), Inches(0.7),
   "Separate actor-critic networks\nLearnable per-dimension log-σ for exploration",
   fs=14, color=MEDIUM_GRAY, ls=1.3)

tb(sl, Inches(0.6), Inches(6.0), Inches(6.0), Inches(0.8),
   "Design rationale: Quaternion avoids gimbal lock; fin angles/rates give "
   "proprioceptive feedback on actuator state (critical given servo lag); "
   "normalized RPM informs policy of available fin authority.",
   fs=13, color=MEDIUM_GRAY, ls=1.25)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Initialization Priors
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Critical: Initialization Priors on Action Distribution", 11)

add_rect(sl, Inches(0.4), Inches(1.15), Inches(6.0), Inches(3.2), LIGHT_BLUE)
tb(sl, Inches(0.7), Inches(1.25), Inches(5.4), Inches(0.35),
   "1. Throttle Bias Initialization", fs=20, color=DARK_BLUE, bold=True)
# Equation
add_eq(sl,
       r"$b_{\rm thr} = \mathrm{atanh}(2\times 0.78 - 1) \approx 1.05$",
       Inches(0.7), Inches(1.7), height=Inches(0.55), fontsize=20)
add_eq(sl,
       r"$u_{\rm hover} = \sqrt{mg / T_{\max}} \approx 0.88$",
       Inches(0.7), Inches(2.35), height=Inches(0.55), fontsize=20)
bullets(sl, Inches(0.7), Inches(3.0), Inches(5.4), Inches(1.2), [
    "Initial mean throttle ≈ 0.78 (near hover)",
    "Without this: default 0.5 → immediate free-fall",
    "Agent converges to zero-throttle local optimum",
], fs=15, spacing=4)

add_rect(sl, Inches(6.8), Inches(1.15), Inches(6.0), Inches(3.2), LIGHT_BLUE)
tb(sl, Inches(7.1), Inches(1.25), Inches(5.4), Inches(0.35),
   "2. Per-Channel Exploration Noise", fs=20, color=DARK_BLUE, bold=True)
add_eq(sl,
       r"$\log\sigma_{\rm fin} = -2.0 \;\Rightarrow\; \sigma \approx 0.14$",
       Inches(7.1), Inches(1.7), height=Inches(0.55), fontsize=20)
add_eq(sl,
       r"$\log\sigma_{\rm thr} = -1.0 \;\Rightarrow\; \sigma \approx 0.37$",
       Inches(7.1), Inches(2.35), height=Inches(0.55), fontsize=20)
bullets(sl, Inches(7.1), Inches(3.0), Inches(5.4), Inches(1.2), [
    "Asymmetric: quieter fins, louder throttle",
    "Prevents lateral oscillations from fin noise",
    "Maintains altitude exploration range",
], fs=15, spacing=4)

add_rect(sl, Inches(0.4), Inches(4.55), Inches(12.4), Inches(2.6), WARM_RED)
tb(sl, Inches(0.7), Inches(4.65), Inches(11.8), Inches(0.35),
   "WITHOUT THESE PRIORS — FAILURE MODES", fs=18, color=RED, bold=True)
tb(sl, Inches(0.7), Inches(5.1), Inches(5.5), Inches(1.7),
   "Default throttle (0.5):\n"
   "→ Immediate free-fall\n"
   "→ Large negative per-step rewards\n"
   "→ Converges to zero-throttle equilibrium\n"
   "→ Agent learns: \"crash fast = minimize cost\"",
   fs=15, color=DARK_GRAY, ls=1.3)
tb(sl, Inches(6.8), Inches(5.1), Inches(5.5), Inches(1.7),
   "Uniform exploration noise:\n"
   "→ Excessive lateral oscillations from fins\n"
   "→ Vehicle destabilized before learning descent\n"
   "→ No stable trajectories to learn from\n"
   "→ Policy never discovers productive control basin",
   fs=15, color=DARK_GRAY, ls=1.3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Reward Design
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Reward Design: Terminal vs. Per-Step Balance", 12)

section_label(sl, Inches(0.6), Inches(1.15), Inches(5.5),
              "LANDING TASK REWARD WEIGHTS")
reward_data = [
    ["Term", "Weight"],
    ["Alive bonus", "+0.10"],
    ["3D position error", "−0.05"],
    ["Horizontal position error", "−0.30"],
    ["Attitude error (tilt)", "−0.20"],
    ["Angular velocity", "−0.10"],
    ["Control effort (fin |α|)", "−0.02"],
    ["Control rate (fin |α̇|)", "−0.01"],
    ["Vertical speed shaping", "−0.30"],
    ["Delta-v cost (T/T_max)²", "−0.20"],
    ["─── Terminal ───", "───"],
    ["Crash penalty", "−200"],
    ["Touchdown softness", "+25·e^(−v)"],
    ["Pad accuracy", "+200·e^(−d)"],
    ["Landing success", "+250"],
]
make_table(sl, Inches(0.5), Inches(1.55), Inches(5.5), Inches(5.2),
           len(reward_data), 2, reward_data,
           col_widths=[Inches(3.5), Inches(2.0)], fs=12)

add_rect(sl, Inches(6.8), Inches(1.15), Inches(6.0), Inches(3.2), WARM_YELLOW)
tb(sl, Inches(7.1), Inches(1.25), Inches(5.4), Inches(0.35),
   "THE MAGNITUDE BALANCE PROBLEM", fs=17, color=ACCENT_ORANGE, bold=True)

# Equation for integrated cost
add_eq(sl,
       r"$\sum_i |w_i|\,\bar{\phi}_i \cdot T \;\sim\; \mathcal{O}(100)$"
       r"$\quad (T=900,\;30\,\mathrm{s})$",
       Inches(7.1), Inches(1.7), height=Inches(0.55), fontsize=18)

tb(sl, Inches(7.1), Inches(2.35), Inches(5.4), Inches(1.8),
   "If terminal rewards < O(100):\n"
   "→ PPO optimizes per-step minimization\n"
   "→ Zero throttle = crash fast = minimize cost\n\n"
   "Solution: terminal magnitudes (≥ 200) must\n"
   "clearly dominate integrated per-step budget",
   fs=15, color=DARK_GRAY, ls=1.3)

add_rect(sl, Inches(6.8), Inches(4.55), Inches(6.0), Inches(2.2), LIGHT_BLUE)
tb(sl, Inches(7.1), Inches(4.65), Inches(5.4), Inches(0.35),
   "Horizontal Guidance Shaping", fs=16, color=DARK_BLUE, bold=True)
add_eq(sl,
       r"$r_{\rm pad} = +200\cdot e^{-d_{\rm horiz}},\quad r_{\rm soft} = +25\cdot e^{-v_{\rm impact}}$",
       Inches(7.1), Inches(5.1), height=Inches(0.55), fontsize=17)
tb(sl, Inches(7.1), Inches(5.75), Inches(5.4), Inches(0.8),
   "Dense lateral gradient (w = −0.30) throughout descent\n"
   "Exponential pad bonus has negligible gradient for d > 2 m",
   fs=14, color=DARK_GRAY, ls=1.3)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Training Procedure & Curriculum
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Training Procedure & Spawn-Position Curriculum", 13)

section_label(sl, Inches(0.6), Inches(1.15), Inches(5.5), "TWO TRAINING TASKS")

add_rect(sl, Inches(0.5), Inches(1.55), Inches(5.8), Inches(2.1), LIGHT_BLUE)
tb(sl, Inches(0.8), Inches(1.65), Inches(5.2), Inches(0.35),
   "Hover Stabilization", fs=18, color=DARK_BLUE, bold=True)
bullets(sl, Inches(0.8), Inches(2.05), Inches(5.2), Inches(1.4), [
    "Residual-PID: PID baseline + PPO correction (scale = 0.05)",
    "2,048 environment steps,  LR = 10⁻⁶",
    "Rollout: 64 steps, 2 minibatches, 2 epochs",
], fs=14, spacing=4)

add_rect(sl, Inches(0.5), Inches(3.9), Inches(5.8), Inches(2.1), WARM_GREEN)
tb(sl, Inches(0.8), Inches(4.0), Inches(5.2), Inches(0.35),
   "Autonomous Landing", fs=18, color=GREEN, bold=True)
bullets(sl, Inches(0.8), Inches(4.4), Inches(5.2), Inches(1.4), [
    "Pure PPO (no PID baseline) — full trajectory learning",
    "5,000,000 environment steps,  LR = 3×10⁻⁴",
    "Rollout: 128 steps, 8 minibatches, 4 epochs",
], fs=14, spacing=4)

section_label(sl, Inches(7.0), Inches(1.15), Inches(5.5), "SPAWN-POSITION CURRICULUM")
add_rect(sl, Inches(6.8), Inches(1.55), Inches(6.0), Inches(5.0), LIGHT_GRAY)

stages = [
    ("Start (0 steps)",   r"$\pm 0.5\,\mathrm{m}$ XY spawn box",   "Learn on-pad touchdown"),
    ("Linear anneal",     r"$\pm 0.5 \to \pm 2.0\,\mathrm{m}$",    "Gradually increase lateral offset"),
    ("End (3M steps)",    r"$\pm 2.0\,\mathrm{m}$ XY (full range)", "Full task difficulty"),
    ("Altitude",          r"$8\text{–}12\,\mathrm{m}$ (constant)",  "Same throughout training"),
    ("Evaluation",        "Full ±2.0 m range",                      "Always un-curricularized"),
]
for i, (stage, value, desc) in enumerate(stages):
    y = Inches(1.75) + Inches(i * 0.9)
    tb(sl, Inches(7.1), y, Inches(2.5), Inches(0.35),
       stage, fs=15, color=DARK_BLUE, bold=True)
    # Render the value as math if it contains $
    if "$" in value:
        add_eq(sl, value, Inches(9.5), y, height=Inches(0.38), fontsize=16)
    else:
        tb(sl, Inches(9.5), y, Inches(3.0), Inches(0.35),
           value, fs=15, color=DARK_GRAY)
    tb(sl, Inches(7.1), y + Inches(0.38), Inches(5.5), Inches(0.35),
       desc, fs=12, color=MEDIUM_GRAY)

tb(sl, Inches(0.6), Inches(6.2), Inches(12.0), Inches(0.6),
   "Rationale: Narrow spawn box → learn soft-touchdown first, then generalize to larger offsets. "
   "Evaluation always uses the full range.",
   fs=14, color=MEDIUM_GRAY, ls=1.25)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Hover Results
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Results: Hover Stabilization", 14)

section_label(sl, Inches(0.6), Inches(1.15), Inches(5.5), "HOVER PERFORMANCE METRICS")
hover_data = [
    ["Metric", "Value", "Threshold"],
    ["Mean position error", "0.085 m", "< 0.5 m  ✓"],
    ["Max position error",  "0.215 m", "< 0.5 m  ✓"],
    ["Mean tilt",           "< 0.001 rad", "—"],
    ["Max tilt",            "< 0.001 rad", "—"],
    ["Mean angular rate",   "1.1×10⁻⁵ rad/s", "—"],
    ["Evaluation",          "PASSED", "—"],
]
make_table(sl, Inches(0.5), Inches(1.55), Inches(6.0), Inches(3.2),
           len(hover_data), 3, hover_data,
           col_widths=[Inches(2.5), Inches(2.0), Inches(1.5)], fs=14)

add_rect(sl, Inches(0.5), Inches(5.0), Inches(6.0), Inches(1.7), LIGHT_BLUE)
tb(sl, Inches(0.8), Inches(5.1), Inches(5.4), Inches(0.35),
   "Key Findings", fs=16, color=DARK_BLUE, bold=True)
bullets(sl, Inches(0.8), Inches(5.5), Inches(5.4), Inches(1.0), [
    "PID provides effective attitude stabilization",
    "RL residual corrects for position drift",
    "Only 2,048 steps needed — very fast convergence",
], fs=14, spacing=4)

for i, (val, label, clr) in enumerate([
    ("0.085 m",    "Mean Position Error", GREEN),
    ("< 0.001 rad","Mean Tilt",           GREEN),
    ("2,048",      "Training Steps",      MEDIUM_BLUE),
]):
    y = Inches(1.3) + Inches(i * 1.8)
    add_rect(sl, Inches(7.0), y, Inches(5.5), Inches(1.5), LIGHT_GRAY)
    tb(sl, Inches(7.3), y + Inches(0.15), Inches(4.9), Inches(0.7),
       val, fs=40, color=clr, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, Inches(7.3), y + Inches(0.9), Inches(4.9), Inches(0.4),
       label, fs=16, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Landing Results
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Results: Autonomous Landing", 15)

section_label(sl, Inches(0.6), Inches(1.15), Inches(6.0),
              "LANDING EVALUATION METRICS (128 PARALLEL EPISODES)")
landing_data = [
    ["Metric", "Value"],
    ["Landing rate",                "100%"],
    ["Mean touchdown speed",        "0.126 m/s"],
    ["Max touchdown speed",         "0.347 m/s"],
    ["Mean pad distance",           "1.89 m"],
    ["Max pad distance",            "4.56 m"],
    ["Mean throttle (eval)",        "0.759"],
    ["Max downward speed",          "9.75 m/s"],
    ["On-pad fraction (d < 0.5 m)", "2%"],
]
make_table(sl, Inches(0.5), Inches(1.55), Inches(5.5), Inches(3.8),
           len(landing_data), 2, landing_data,
           col_widths=[Inches(3.5), Inches(2.0)], fs=14)

for i, (val, label, clr) in enumerate([
    ("100%",      "Landing Rate",          GREEN),
    ("0.126 m/s", "Mean Touchdown Speed",  GREEN),
    ("1.89 m",    "Mean Pad Distance",     ACCENT_ORANGE),
    ("5M",        "Training Steps",        MEDIUM_BLUE),
]):
    x = Inches(7.0) + Inches((i % 2) * 3.0)
    y = Inches(1.3) + Inches((i // 2) * 2.0)
    add_rect(sl, x, y, Inches(2.7), Inches(1.7), LIGHT_GRAY)
    tb(sl, x + Inches(0.1), y + Inches(0.15), Inches(2.5), Inches(0.7),
       val, fs=32, color=clr, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), y + Inches(1.0), Inches(2.5), Inches(0.5),
       label, fs=13, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2), WARM_YELLOW)
tb(sl, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.9),
   "✓  100% landing rate with soft touchdowns — the policy reliably descends and lands safely\n"
   "⚠  Lateral accuracy limited: 1.89 m mean pad distance, only 2% on-pad — "
   "lateral guidance to the pad center remains an open challenge",
   fs=16, color=DARK_GRAY, ls=1.35)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — Training Curves
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Training Progression & Analysis", 16)

tc_path = os.path.join(FIGURES_DIR, "training_curves.png")
if os.path.exists(tc_path):
    img(sl, tc_path, Inches(0.5), Inches(1.15), h=Inches(4.2))
else:
    add_rect(sl, Inches(0.5), Inches(1.15), Inches(6.5), Inches(4.0), LIGHT_GRAY)
    tb(sl, Inches(1.5), Inches(2.7), Inches(4.5), Inches(0.5),
       "[training_curves.png\nMean Episode Reward over 5M Steps]",
       fs=14, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

section_label(sl, Inches(7.5), Inches(1.15), Inches(5.5), "TRAINING DYNAMICS")
bullets(sl, Inches(7.5), Inches(1.55), Inches(5.3), Inches(4.0), [
    "0–1M steps: reward ≈ −3 (per-step costs dominate)",
    "1–2M steps: transition to positive rewards",
    "2–5M steps: steady improvement to ≈ +1.75",
    "Monotonic upward trend confirms learning",
    "100% landing rate achieved early, maintained throughout",
], fs=16, spacing=8)

add_rect(sl, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.4), LIGHT_BLUE)
tb(sl, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.35),
   "Interpreting the Reward Crossover at ~2M Steps", fs=16, color=DARK_BLUE, bold=True)
tb(sl, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.7),
   "Negative → Positive: integrated per-step costs give way to terminal landing bonuses "
   "(+250 success + up to +200 pad accuracy). Marks when the majority of episodes end "
   "in successful landings rather than timeouts or crashes.",
   fs=15, color=DARK_GRAY, ls=1.25)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — Lateral Accuracy Challenge
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Discussion: The Lateral Accuracy Challenge", 17)

challenges = [
    ("Throttle–Fin Coupling",
     r"Lateral authority $\propto$ throttle$^2$. Reducing throttle to descend "
     r"simultaneously reduces fin effectiveness — fundamental tension between "
     r"descent rate and lateral controllability."),
    ("Actuator Bandwidth",
     r"MG996R servos ($\tau = 0.05\,\mathrm{s}$, $\dot{\alpha}_{\max} = 7.54\,\mathrm{rad/s}$) "
     r"limit correction speed, especially during final approach."),
    ("Reward Gradient at Distance",
     r"Exponential pad bonus $e^{-d}$ has negligible gradient for $d > 2\,\mathrm{m}$. "
     r"Linear horizontal penalty ($w = -0.30$) may be insufficient vs. vertical/delta-v terms."),
    ("Curriculum Interaction",
     r"Narrow initial spawn helps learn touchdown but may not provide sufficient "
     r"training signal for large lateral corrections from the full $\pm 2.0\,\mathrm{m}$ range."),
]

for i, (title, desc) in enumerate(challenges):
    y = Inches(1.2) + Inches(i * 1.45)
    oval_num(sl, Inches(0.6), y + Inches(0.05), str(i + 1), bg=ACCENT_ORANGE)
    tb(sl, Inches(1.3), y, Inches(11.5), Inches(0.35),
       title, fs=18, color=DARK_BLUE, bold=True)
    # Render description — strip $ for plain text version, render key eq separately
    plain = (desc.replace(r"$\propto$", "∝")
                 .replace(r"$^2$", "²")
                 .replace(r"$\tau = 0.05\,\mathrm{s}$", "τ = 0.05 s")
                 .replace(r"$\dot{\alpha}_{\max} = 7.54\,\mathrm{rad/s}$", "α̇_max = 7.54 rad/s")
                 .replace(r"$e^{-d}$", "e^(−d)")
                 .replace(r"$d > 2\,\mathrm{m}$", "d > 2 m")
                 .replace(r"$w = -0.30$", "w = −0.30")
                 .replace(r"$\pm 2.0\,\mathrm{m}$", "±2.0 m"))
    tb(sl, Inches(1.3), y + Inches(0.4), Inches(11.5), Inches(0.8),
       plain, fs=14, color=DARK_GRAY, ls=1.25)

add_rect(sl, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), LIGHT_BLUE)
tb(sl, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
   "vs. Quadrotors: Independent, throttle-decoupled attitude authority. "
   "TVC's coupled control means standard quadrotor RL techniques don't directly transfer.",
   fs=14, color=DARK_BLUE, ls=1.2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — Sim-to-Real
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Sim-to-Real Transfer Considerations", 18)

gaps = [
    ("Aerodynamic Model Fidelity",
     "Flat-plate fin model may miss flow separation, wake interactions, and ground-effect changes.",
     "CFD validation, wind tunnel data"),
    ("Servo Nonlinearities",
     "Real MG996R servos exhibit backlash, load-dependent speed variation, and temperature drift.",
     "System identification on hardware"),
    ("Battery Voltage Sag",
     "6S LiPo voltage drops under load, reducing available thrust — not currently modeled.",
     "Voltage-dependent thrust model"),
    ("Sensor Noise",
     "Simulation provides clean state observations; physical testbed uses noisy IMU and position estimates.",
     "Observation noise injection during training"),
]
for i, (title, desc, mitigation) in enumerate(gaps):
    y = Inches(1.2) + Inches(i * 1.4)
    add_rect(sl, Inches(0.5), y, Inches(8.0), Inches(1.15),
             LIGHT_GRAY if i % 2 == 0 else WHITE)
    tb(sl, Inches(0.8), y + Inches(0.05), Inches(7.5), Inches(0.3),
       title, fs=16, color=DARK_BLUE, bold=True)
    tb(sl, Inches(0.8), y + Inches(0.35), Inches(7.5), Inches(0.6),
       desc, fs=13, color=DARK_GRAY, ls=1.2)
    add_rect(sl, Inches(8.8), y, Inches(4.0), Inches(1.15), WARM_GREEN)
    tb(sl, Inches(9.0), y + Inches(0.05), Inches(3.6), Inches(0.25),
       "Mitigation:", fs=12, color=GREEN, bold=True)
    tb(sl, Inches(9.0), y + Inches(0.35), Inches(3.6), Inches(0.6),
       mitigation, fs=13, color=DARK_GRAY, ls=1.2)

add_rect(sl, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8), LIGHT_BLUE)
tb(sl, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.5),
   "Primary strategy: Domain randomization of aerodynamic coefficients, servo parameters, "
   "mass properties, and sensor noise during training to improve transfer robustness.",
   fs=15, color=DARK_BLUE, ls=1.2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — Future Work
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
header(sl, "Future Work", 19)

future = [
    ("GTrXL Policy",
     "Replace feed-forward MLP with Gated Transformer-XL for temporal context — "
     "handle wind disturbances and improve trajectory planning. "
     "Recent aerospace results (Federici et al., Carradori et al.) show promise.",
     MEDIUM_BLUE),
    ("Domain Randomization",
     "Randomize aerodynamic coefficients, servo parameters, mass properties, "
     "and sensor noise during training for robust sim-to-real transfer.",
     GREEN),
    ("Lateral Guidance",
     "Potential-based reward shaping (Ng et al., 1999), hierarchical policies "
     "(separate lateral/vertical controllers), extended curriculum strategies.",
     ACCENT_ORANGE),
    ("Hardware Validation",
     "Deploy trained policies on the physical EDF drone testbed. "
     "Validate sim-to-real transfer and characterize the reality gap.",
     RED),
    ("Wind Disturbance Rejection",
     "Train and evaluate under stochastic wind conditions using the simulation's "
     "wind model with 27-dim observation space including wind estimates.",
     RGBColor(0x7B, 0x1F, 0xA2)),
]
for i, (title, desc, color) in enumerate(future):
    y = Inches(1.2) + Inches(i * 1.15)
    add_rect(sl, Inches(0.5), y, Inches(0.12), Inches(0.9), color)
    tb(sl, Inches(0.9), y, Inches(4.0), Inches(0.35),
       title, fs=18, color=color, bold=True)
    tb(sl, Inches(0.9), y + Inches(0.35), Inches(11.5), Inches(0.55),
       desc, fs=14, color=DARK_GRAY, ls=1.2)

add_rect(sl, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7), LIGHT_BLUE)
tb(sl, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
   "Near-term priority: GTrXL integration + domain randomization → hardware flight tests. "
   "The 24-dim observation space is already compatible with transformer sequence input.",
   fs=15, color=DARK_BLUE, ls=1.2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(2.5), DARK_BLUE)
add_rect(sl, Inches(0), Inches(2.5), SLIDE_W, Inches(0.08), ACCENT_ORANGE)

tb(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.6),
   "Conclusion", fs=34, color=WHITE, bold=True)
tb(sl, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.2),
   "PPO can learn thrust-vectored fin control for an EDF VTOL vehicle in simulation — "
   "achieving 100% landing rate with soft touchdowns, opening the path to sim-to-real transfer.",
   fs=20, color=RGBColor(0xBB, 0xD5, 0xED), ls=1.3)

section_label(sl, Inches(0.6), Inches(2.8), Inches(12.0), "KEY CONTRIBUTIONS")
contributions = [
    "GPU-accelerated Isaac Sim environment with semi-empirical aerodynamic, propulsion, and actuator models",
    "PPO hover policy: 0.085 m mean error in only 2,048 steps (residual-PID architecture)",
    "PPO landing policy: 100% success rate, 0.126 m/s mean touchdown speed over 5M steps",
    "Critical insight: initialization priors on action distribution prevent degenerate local optima",
    "Reward design principle: terminal magnitudes must dominate integrated per-step costs",
]
for i, contrib in enumerate(contributions):
    y = Inches(3.2) + Inches(i * 0.55)
    tb(sl, Inches(0.6), y, Inches(0.4), Inches(0.4), "✓", fs=18, color=GREEN, bold=True)
    tb(sl, Inches(1.0), y, Inches(11.5), Inches(0.45), contrib, fs=16, color=DARK_GRAY, ls=1.15)

add_rect(sl, Inches(0), Inches(6.2), SLIDE_W, Inches(1.3), LIGHT_GRAY)
tb(sl, Inches(0.8), Inches(6.3), Inches(6), Inches(0.4),
   "Tang Zijian (Jacob Tang)  |  Minnesota State University, Mankato",
   fs=14, color=DARK_GRAY)
tb(sl, Inches(0.8), Inches(6.7), Inches(10), Inches(0.4),
   "Code & Data: github.com/Jacob19999/Transformer-rl-retro-propulsion",
   fs=14, color=MEDIUM_BLUE, bold=True)
tb(sl, Inches(9.0), Inches(6.3), Inches(4.0), Inches(0.8),
   "Thank You!\nQuestions?",
   fs=28, color=DARK_BLUE, bold=True, align=PP_ALIGN.RIGHT)

footer(sl); slide_num(sl, 20)

# ═══════════════════════════════════════════════════════════════
# ADD SPEAKER NOTES TO ALL SLIDES
# ═══════════════════════════════════════════════════════════════
from lxml import etree

def set_notes(slide, text):
    """Set the speaker notes for a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

notes_text = [
    # Slide 1
    ("Good morning/afternoon everyone. My name is Jacob Tang from Minnesota State University, Mankato. "
     "Today I'll be presenting our work on using reinforcement learning — specifically Proximal Policy Optimization — "
     "to control a thrust-vectored electric ducted fan drone in simulation. "
     "This is a vehicle that steers not by spinning multiple rotors at different speeds like a typical quadcopter, "
     "but by deflecting small fins inside its own exhaust stream. "
     "It's a fundamentally different and much harder control problem, and we'll show how deep RL can tackle it.\n\n"
     "[~30 seconds]"),
    # Slide 2
    ("Here's our roadmap for the next 20 minutes. We'll start with why this problem matters and what makes it hard. "
     "Then we'll walk through the vehicle itself, the simulation environment we built in NVIDIA Isaac Sim, "
     "the physics models under the hood, our PPO training setup and reward design, and then the results. "
     "The key question: can deep RL learn to control a vehicle that steers by deflecting fins in its own exhaust stream? "
     "Spoiler — yes, but with some interesting caveats.\n\n[~30 seconds]"),
    # Slide 3
    ("So why are we doing this? Most drones today are multirotors — they control attitude by spinning each rotor at a different speed. "
     "Our vehicle uses thrust vector control: one electric ducted fan, four flat-plate fins in the exhaust stream.\n\n"
     "The fundamental challenge is the coupling shown at the bottom right: fin force is proportional to throttle squared. "
     "At 50% throttle, you only have 25% of your maximum control authority. "
     "This is a deeply nonlinear, tightly coupled control problem where RL shines — learning the full nonlinear mapping end-to-end.\n\n"
     "[~1.5 minutes]"),
    # Slide 4
    ("Briefly on related work. RL for aerial vehicles has made great progress — Hwangbo showed PPO can learn agile quadrotor control, "
     "Koch compared multiple RL algorithms for attitude control, and GPU-parallel simulation from Isaac Gym has been a game-changer.\n\n"
     "On the TVC side, there's a long history in rocketry, but applying RL to jet-vane TVC for VTOL vehicles is largely unexplored. "
     "The research gap is clear: RL for thrust-vectored fin control of EDF vehicles, with the unique challenges of coupled "
     "throttle-authority, hobby-grade actuators, and gyroscopic effects.\n\n[~1 minute]"),
    # Slide 5
    ("Here's our testbed. 3.1 kilogram EDF drone, single FMS 90mm 12-blade ducted fan, 6S LiPo battery. "
     "Four MG996R hobby servos actuate flat-plate jet vanes in the exhaust stream, arranged in a cruciform pattern.\n\n"
     "The control architecture gives us full 6-DOF authority through just 5 actuator commands: "
     "four fin angles for pitch, roll, and yaw, plus one throttle for altitude.\n\n[~1 minute]"),
    # Slide 6
    ("We built our simulation on NVIDIA Isaac Sim using the Isaac Lab framework. "
     "The key enabler is GPU-accelerated parallel simulation — 128 copies of the drone simultaneously on a single GPU. "
     "Physics runs at 120 Hz, RL agent makes decisions at 30 Hz with a decimation factor of 4.\n\n[~1 minute]"),
    # Slide 7
    ("The EDF propulsion model has two key components: first-order spool dynamics with a 0.15 second time constant, "
     "and quadratic thrust scaling with rotor speed squared.\n\n"
     "The three torque components are what make this interesting. Gyroscopic precession couples body rotation to rotor "
     "angular momentum — if the vehicle pitches while the rotor spins at 4300 rad/s, you get a yaw torque. "
     "These real physical effects must be captured in simulation.\n\n[~1 minute]"),
    # Slide 8
    ("The fin aerodynamic model is semi-empirical. Each fin sees a dynamic pressure that depends on exhaust velocity "
     "and throttle squared — there's that coupling again.\n\n"
     "The critical point: this quadratic coupling between throttle and fin effectiveness is THE defining characteristic "
     "of jet-vane TVC. When our vehicle descends and reduces throttle, it simultaneously loses its ability to correct "
     "lateral errors. This is the core tension the RL agent must learn to navigate.\n\n[~1 minute]"),
    # Slide 9
    ("We use standard PPO with clipped surrogate objectives. PPO is an on-policy actor-critic algorithm that constrains "
     "how much the policy can change in each update — preventing catastrophically large gradient steps while maintaining "
     "good sample efficiency. It's become the workhorse of continuous control RL.\n\n"
     "Our setup uses GAE for advantage estimation with gamma 0.99 and lambda 0.95, plus adaptive KL early stopping.\n\n[~1 minute]"),
    # Slide 10
    ("The observation space is 24-dimensional. Key design choices: we include fin angles and rates as proprioceptive feedback "
     "because the servo lag is significant — the policy needs to know where the fins actually are. "
     "And we include normalized motor RPM so the policy knows how much fin authority it currently has.\n\n"
     "The action space is just 5 dimensions — four fin deflections and one throttle. "
     "The network is straightforward: separate actor and critic, each with two hidden layers of 256 units.\n\n[~1 minute]"),
    # Slide 11
    ("This slide is arguably the most important practical insight from our work. Two initialization choices are critical.\n\n"
     "First, throttle bias initialization. We set the actor's output bias so the initial mean throttle maps to about 0.78 — "
     "near the hover throttle of 0.88. Without this, the default 0.5 means immediate free-fall, and PPO converges to a "
     "zero-throttle equilibrium — the agent learns that crashing fast minimizes cost.\n\n"
     "Second, per-channel exploration noise. Quieter fins, louder throttle. With uniform noise, the fins oscillate wildly, "
     "destabilizing the vehicle before it can learn anything useful.\n\n"
     "These are physically motivated initialization priors, not ad hoc hacks. Without them, training simply fails.\n\n[~1.5 minutes]"),
    # Slide 12
    ("Reward design for landing is tricky, and we learned an important lesson about magnitude balance.\n\n"
     "For a 30-second episode at 30 Hz, that's 900 steps. The integrated per-step costs can easily reach order 100. "
     "If your terminal rewards aren't substantially larger, PPO will optimize per-step cost minimization instead of actually landing. "
     "We set terminal magnitudes at 200 or above to clearly dominate the per-step budget.\n\n"
     "The horizontal guidance shaping at weight -0.30 provides a dense lateral gradient throughout descent.\n\n[~1.5 minutes]"),
    # Slide 13
    ("We train two tasks. Hover uses a residual-PID architecture — a PID controller provides the baseline, "
     "and PPO learns a small bounded correction. This converges in just 2,048 steps.\n\n"
     "Landing is the main event: pure PPO, no PID baseline, 5 million steps. "
     "The spawn-position curriculum starts with a narrow ±0.5 meter spawn box, linearly annealing to ±2.0 meters. "
     "Evaluation always uses the full range.\n\n[~1 minute]"),
    # Slide 14
    ("Hover results are clean. Mean position error of 0.085 meters, well within the 0.5 meter threshold. "
     "Near-zero tilt and angular rates. The remarkable thing is this only took 2,048 environment steps to converge. "
     "The residual architecture is very efficient when you have a reasonable baseline controller.\n\n[~45 seconds]"),
    # Slide 15
    ("Landing results are more nuanced. The headline: 100% landing rate with a mean touchdown speed of 0.126 m/s — "
     "a very soft landing, well below the 0.5 m/s contact detection threshold.\n\n"
     "However, lateral accuracy is the open challenge. Mean pad distance of 1.89 meters, only 2% on-pad. "
     "The policy reliably descends and lands softly, but doesn't consistently steer toward the pad center. "
     "It's learned 'land safely anywhere' but not 'land precisely here'.\n\n[~1.5 minutes]"),
    # Slide 16
    ("The training curve tells a clear story. Reward starts around -3 — per-step costs dominate. "
     "Around 2 million steps, we see the crossover to positive rewards, meaning terminal landing bonuses "
     "now outweigh the per-step costs. By 5 million steps, mean reward reaches about +1.75.\n\n"
     "That negative-to-positive crossover marks when the majority of episodes end in successful landings "
     "rather than timeouts or crashes.\n\n[~45 seconds]"),
    # Slide 17
    ("Let me dig into why lateral accuracy is hard. Four factors contribute.\n\n"
     "First, the throttle-fin coupling — descending means less lateral authority. "
     "Second, servo bandwidth limits correction rate during final approach. "
     "Third, the exponential pad bonus has essentially zero gradient beyond 2 meters. "
     "Fourth, the curriculum may not provide enough training signal for large lateral corrections.\n\n"
     "Compared to quadrotor RL, this is a fundamentally harder lateral control problem because "
     "attitude authority is coupled to thrust state.\n\n[~1 minute]"),
    # Slide 18
    ("For sim-to-real transfer, we've identified four main gap sources: aerodynamic model fidelity, "
     "servo nonlinearities, battery voltage sag, and sensor noise. "
     "Each has a mitigation path, and the overarching strategy is domain randomization — "
     "randomizing these parameters during training so the policy learns to be robust to the uncertainty.\n\n[~1 minute]"),
    # Slide 19
    ("Five directions going forward. First, replacing the feed-forward MLP with a Gated Transformer-XL architecture "
     "to give the policy temporal context — this should help with wind disturbance rejection and trajectory planning. "
     "Second, domain randomization for sim-to-real robustness. Third, improving lateral guidance. "
     "Fourth, the ultimate goal — deploying on the physical drone testbed. "
     "Fifth, training under stochastic wind conditions.\n\n[~1 minute]"),
    # Slide 20
    ("To summarize: we've shown that PPO can learn thrust-vectored fin control for an EDF VTOL vehicle in simulation. "
     "The hover policy achieves 8.5 centimeter accuracy in just 2,048 steps. "
     "The landing policy achieves 100% success rate with soft touchdowns over 5 million steps. "
     "We've identified critical initialization priors and reward design principles specific to TVC vehicles. "
     "The code and data are publicly available on GitHub. Thank you — I'm happy to take questions.\n\n[~30 seconds]"),
]

for i, (slide, note) in enumerate(zip(prs.slides, notes_text)):
    set_notes(slide, note)

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "TVC_RL_Presentation_v2.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
